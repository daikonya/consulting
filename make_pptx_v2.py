# -*- coding: utf-8 -*-
"""
提案書 v2.pptx  —  Brain x TOJISM  (revised)
Slides:
  01 表紙
  02 アジェンダ
  03 ブランディング戦略（事業承継ターゲット追加）
  04 集客ファネル（KPI revised）
  05 サービスプラン比較（ライト/スタンダード/プレミアム）
  06 サービス詳細①  LP（別途投影）
  07 サービス詳細②  公式LINE構築 + LINE広告
  08 サービス詳細③  ブランディング戦略支援
  09 サービス詳細④  月次マーケ支援
  10 サービス詳細⑤  SNS動画制作・運用
  11 サービス詳細⑥  AI活用・自動ブログ
"""

import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ─────────────────────────────────────────────
DARK   = RGBColor(0x1C, 0x1A, 0x16)
DARK2  = RGBColor(0x26, 0x23, 0x1E)
DARK3  = RGBColor(0x22, 0x20, 0x1C)
GOLD   = RGBColor(0xB8, 0x93, 0x5A)
GOLD_L = RGBColor(0xD4, 0xAF, 0x7A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CREAM  = RGBColor(0xF5, 0xF2, 0xED)
LGRAY  = RGBColor(0xEE, 0xEA, 0xE3)
MGRAY  = RGBColor(0x6B, 0x63, 0x59)
BGRAY  = RGBColor(0xDD, 0xD5, 0xC5)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
BLUE   = RGBColor(0x2E, 0x86, 0xC1)
TEAL   = RGBColor(0x1A, 0x7A, 0x5E)
LINE_GREEN = RGBColor(0x06, 0xC7, 0x55)

SW = Inches(13.33)
SH = Inches(7.5)
JP  = "Yu Gothic"
JPS = "Yu Mincho"


# ── Primitive helpers ────────────────────────────────────

def new_prs():
    p = Presentation()
    p.slide_width  = SW
    p.slide_height = SH
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=DARK):
    r = slide.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()

def box(slide, l, t, w, h, fill=None, lc=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if lc:   s.line.color.rgb = lc; s.line.width = lw
    else:    s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h,
       fn=JP, sz=13, fc=WHITE, bold=False,
       align=PP_ALIGN.LEFT, wrap=True):
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = fn
    run.font.size  = Pt(sz)
    run.font.color.rgb = fc
    run.font.bold  = bold
    return b

def para(tf, text, fn=JP, sz=12, fc=WHITE, bold=False,
         align=PP_ALIGN.LEFT, sb=Pt(4)):
    p = tf.add_paragraph()
    p.alignment = align; p.space_before = sb
    r = p.add_run()
    r.text = text; r.font.name = fn; r.font.size = Pt(sz)
    r.font.color.rgb = fc; r.font.bold = bold

def accent_bar(slide, color=GOLD):
    box(slide, 0, 0, Inches(0.12), SH, fill=color)

def gold_line(slide, l, t, w=Inches(0.32)):
    box(slide, l, t, w, Inches(0.025), fill=GOLD)

def sec_label(slide, text, l=Inches(0.35), t=Inches(0.38), color=GOLD):
    tb(slide, text, l, t, Inches(10), Inches(0.3), fn=JP, sz=9, fc=color)

def heading(slide, text, l=Inches(0.35), t=Inches(0.7), w=Inches(11), sz=26, fc=WHITE):
    tb(slide, text, l, t, w, Inches(0.9), fn=JPS, sz=sz, fc=fc)

def chip_badge(slide, text, l, t, fill=DARK3, fc=GOLD_L, w=Inches(2.0), h=Inches(0.36)):
    box(slide, l, t, w, h, fill=fill)
    tb(slide, text, l+Inches(0.1), t+Inches(0.05), w-Inches(0.2), h-Inches(0.05),
       fn=JP, sz=10, fc=fc, align=PP_ALIGN.CENTER)

def bullet_list(slide, items, l, t, w, spacing=Inches(0.52),
                dot_color=GOLD, fc=WHITE, sz=12, bg_color=None, row_h=None):
    rh = row_h or spacing
    for i, item in enumerate(items):
        iy = t + i * spacing
        if bg_color:
            box(slide, l, iy, w, rh - Inches(0.04), fill=bg_color)
        box(slide, l, iy + (rh/2) - Inches(0.05), Inches(0.06), Inches(0.06), fill=dot_color)
        tb(slide, item, l+Inches(0.18), iy+Inches(0.06), w-Inches(0.22), rh-Inches(0.08),
           fn=JP, sz=sz, fc=fc, wrap=True)


# ════════════════════════════════════════════════════════
#  SLIDES
# ════════════════════════════════════════════════════════

def s01_cover(prs):
    s = blank(prs)
    bg(s, DARK)
    # left panel
    box(s, 0, 0, Inches(4.3), SH, fill=DARK2)
    box(s, Inches(4.3), 0, Inches(0.06), SH, fill=GOLD)
    # tag
    tb(s, "Business Proposal  /  2026.05.07",
       Inches(4.6), Inches(0.7), Inches(8), Inches(0.35),
       fn=JP, sz=10, fc=RGBColor(0x88,0x80,0x70))
    # main title
    b = s.shapes.add_textbox(Inches(4.6), Inches(1.3), Inches(8.4), Inches(3.0))
    tf = b.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "旅館・飲食業\nコンサルティング\nブランディング戦略"
    r.font.name = JPS; r.font.size = Pt(34); r.font.color.rgb = WHITE
    # sub
    tb(s, "Web戦略 / 集客設計 / サポートプランのご提案",
       Inches(4.6), Inches(4.4), Inches(8), Inches(0.45),
       fn=JP, sz=13, fc=RGBColor(0xAA,0xA0,0x90))
    # left: 提案先
    tb(s, "提案先", Inches(0.35), Inches(1.8), Inches(3.5), Inches(0.3),
       fn=JP, sz=9, fc=GOLD)
    tb(s, "小野 順平 氏\n株式会社TOJISM",
       Inches(0.35), Inches(2.15), Inches(3.5), Inches(0.9),
       fn=JPS, sz=16, fc=WHITE)
    tb(s, "提案者", Inches(0.35), Inches(3.3), Inches(3.5), Inches(0.3),
       fn=JP, sz=9, fc=GOLD)
    tb(s, "合同会社Brain\n鵜木 芳文",
       Inches(0.35), Inches(3.65), Inches(3.5), Inches(0.7),
       fn=JP, sz=13, fc=WHITE)
    # goal box
    box(s, Inches(0.25), Inches(5.3), Inches(3.7), Inches(1.85),
        fill=RGBColor(0x28,0x26,0x22), lc=GOLD, lw=Pt(0.75))
    tb(s, "本日のゴール", Inches(0.45), Inches(5.45), Inches(3.2), Inches(0.3),
       fn=JP, sz=9, fc=GOLD)
    tb(s, "コンサルブランドの方向性に合意し\n具体的な支援プランを決定する",
       Inches(0.45), Inches(5.82), Inches(3.3), Inches(1.1),
       fn=JP, sz=12, fc=WHITE, wrap=True)


def s02_agenda(prs):
    s = blank(prs)
    bg(s, CREAM)
    accent_bar(s)
    sec_label(s, "AGENDA", color=GOLD)
    heading(s, "本日のアジェンダ", fc=DARK, sz=26)
    gold_line(s, Inches(0.35), Inches(1.38))

    rows = [
        ("01", "ブランディング戦略",         "ポジショニング・ターゲット・差別化軸",        "15 min"),
        ("02", "Web戦略・集客ファネル",        "LP / 公式LINE / マーケ計画",                  "15 min"),
        ("03", "サービスプラン比較",          "ライト / スタンダード / プレミアム",           "10 min"),
        ("04", "サービス詳細",               "LP / LINE / ブランディング / マーケ / SNS / AI", "30 min"),
    ]
    for i, (num, title, sub, dur) in enumerate(rows):
        y = Inches(1.6) + i * Inches(1.3)
        box(s, Inches(0.25), y, Inches(12.83), Inches(1.18),
            fill=WHITE, lc=BGRAY, lw=Pt(0.5))
        tb(s, num, Inches(0.45), y+Inches(0.15), Inches(0.7), Inches(0.85),
           fn=JPS, sz=22, fc=GOLD)
        tb(s, title, Inches(1.3), y+Inches(0.13), Inches(8.5), Inches(0.45),
           fn=JP, sz=15, fc=DARK, bold=True)
        tb(s, sub, Inches(1.3), y+Inches(0.6), Inches(8.5), Inches(0.42),
           fn=JP, sz=11, fc=MGRAY)
        box(s, Inches(11.2), y+Inches(0.38), Inches(1.6), Inches(0.36),
            fill=RGBColor(0x99,0x90,0x88))
        tb(s, dur, Inches(11.2), y+Inches(0.45), Inches(1.6), Inches(0.28),
           fn=JP, sz=10, fc=WHITE, align=PP_ALIGN.CENTER)


def s03_branding(prs):
    s = blank(prs)
    bg(s, DARK)
    accent_bar(s)
    sec_label(s, "01 -- BRANDING STRATEGY")
    heading(s, "ターゲット戦略：2つの顧客像")
    gold_line(s, Inches(0.35), Inches(1.38))

    # Target A
    box(s, Inches(0.25), Inches(1.55), Inches(6.2), Inches(4.5),
        fill=DARK3, lc=GOLD, lw=Pt(1.5))
    tb(s, "TARGET A", Inches(0.45), Inches(1.72), Inches(5.8), Inches(0.3),
       fn=JP, sz=9, fc=GOLD)
    tb(s, "経営改善・集客強化を\n求めるオーナー",
       Inches(0.45), Inches(2.05), Inches(5.8), Inches(0.9),
       fn=JPS, sz=18, fc=WHITE)
    items_a = [
        "売上・集客に課題を感じている",
        "予約サイト依存から脱却したい",
        "AIやDXを導入したいが方法がわからない",
        "人材採用・定着に悩んでいる",
    ]
    for j, item in enumerate(items_a):
        iy = Inches(3.1) + j * Inches(0.55)
        box(s, Inches(0.45), iy+Inches(0.14), Inches(0.07), Inches(0.07), fill=GOLD)
        tb(s, item, Inches(0.65), iy+Inches(0.07), Inches(5.6), Inches(0.45),
           fn=JP, sz=12, fc=RGBColor(0xCC,0xC5,0xB8))

    # Target B
    box(s, Inches(6.7), Inches(1.55), Inches(6.38), Inches(4.5),
        fill=RGBColor(0x18,0x2A,0x22), lc=TEAL, lw=Pt(1.5))
    box(s, Inches(6.7), Inches(1.55), Inches(1.5), Inches(0.45), fill=TEAL)
    tb(s, "NEW TARGET", Inches(6.7), Inches(1.62), Inches(1.5), Inches(0.32),
       fn=JP, sz=9, fc=DARK, bold=True, align=PP_ALIGN.CENTER)
    tb(s, "TARGET B", Inches(8.35), Inches(1.72), Inches(4.5), Inches(0.3),
       fn=JP, sz=9, fc=TEAL)
    tb(s, "事業承継・出口戦略を\n検討している経営者",
       Inches(6.9), Inches(2.05), Inches(5.9), Inches(0.9),
       fn=JPS, sz=18, fc=WHITE)
    items_b = [
        "後継者がおらず廃業を検討している",
        "M&A・第三者承継に興味があるが手続きがわからない",
        "引退前に施設の価値を最大化したい",
        "従業員の雇用を守りながら事業を継続させたい",
    ]
    for j, item in enumerate(items_b):
        iy = Inches(3.1) + j * Inches(0.55)
        box(s, Inches(6.9), iy+Inches(0.14), Inches(0.07), Inches(0.07), fill=TEAL)
        tb(s, item, Inches(7.1), iy+Inches(0.07), Inches(5.7), Inches(0.45),
           fn=JP, sz=12, fc=RGBColor(0xCC,0xC5,0xB8))

    # Market size bar
    box(s, Inches(0.25), Inches(6.18), Inches(12.83), Inches(1.1),
        fill=RGBColor(0x26,0x23,0x1E))
    tb(s, "市場規模",
       Inches(0.45), Inches(6.28), Inches(2), Inches(0.3),
       fn=JP, sz=9, fc=GOLD)
    stats = [
        ("飲食店", "約60万店"),
        ("旅館・ホテル", "約5万施設"),
        ("経営者60歳以上", "飲食業 約55%"),
        ("事業承継課題あり（推計）", "約25万事業者"),
    ]
    for i, (lbl, val) in enumerate(stats):
        sx = Inches(0.45) + i * Inches(3.15)
        tb(s, lbl, sx, Inches(6.6), Inches(3), Inches(0.28),
           fn=JP, sz=9, fc=RGBColor(0x88,0x80,0x70))
        tb(s, val, sx, Inches(6.88), Inches(3), Inches(0.28),
           fn=JP, sz=13, fc=WHITE, bold=True)

    # VS arrow
    tb(s, "×", Inches(6.18), Inches(3.45), Inches(0.5), Inches(0.5),
       fn=JP, sz=22, fc=RGBColor(0x55,0x50,0x48), align=PP_ALIGN.CENTER)
    tb(s, "対象者が広がることで\n専門性も高まる",
       Inches(5.7), Inches(5.8), Inches(1.6), Inches(0.6),
       fn=JP, sz=9, fc=RGBColor(0x88,0x80,0x70), align=PP_ALIGN.CENTER)


def s04_funnel(prs):
    s = blank(prs)
    bg(s, DARK)
    accent_bar(s)
    sec_label(s, "02 -- WEB STRATEGY / FUNNEL")
    heading(s, "集客ファネル設計")
    gold_line(s, Inches(0.35), Inches(1.38))

    # Funnel steps
    steps = [
        ("認知",           "Google検索 / セミナー登壇 / 紹介・口コミ / Instagram",   DARK3,   False),
        ("LP",             "信頼構築・サービス訴求・無料相談CTA",                     RGBColor(0x28,0x26,0x22), True),
        ("公式LINE 登録",  "自動フォロー配信・セミナー案内・相談予約導線",            RGBColor(0x06,0x50,0x28), True),
        ("無料相談（60分）","課題ヒアリング → プラン提案",                            DARK3,   False),
        ("契約・支援スタート","ライト / スタンダード / プレミアム",                    DARK3,   False),
    ]
    bw = Inches(6.0); bh = Inches(0.82); bx = Inches(0.25)
    for i, (title, sub, fbg, hl) in enumerate(steps):
        by = Inches(1.55) + i * (bh + Inches(0.18))
        lc = GOLD if hl else RGBColor(0x44,0x40,0x38)
        lw = Pt(1.2) if hl else Pt(0.4)
        box(s, bx, by, bw, bh, fill=fbg, lc=lc, lw=lw)
        tc = GOLD_L if hl else RGBColor(0xCC,0xC5,0xB8)
        tb(s, title, bx+Inches(0.2), by+Inches(0.1), bw-Inches(0.3), Inches(0.38),
           fn=JP, sz=14, fc=tc, bold=hl)
        tb(s, sub, bx+Inches(0.2), by+Inches(0.48), bw-Inches(0.3), Inches(0.3),
           fn=JP, sz=10, fc=RGBColor(0x88,0x80,0x70))
        if i < len(steps)-1:
            tb(s, "v", bx+Inches(2.7), by+bh, Inches(0.6), Inches(0.22),
               fn=JP, sz=11, fc=GOLD, align=PP_ALIGN.CENTER)

    # Right: KPI & Notes
    rx = Inches(6.65)
    # Market note
    box(s, rx, Inches(1.55), Inches(6.43), Inches(1.3), fill=DARK3)
    tb(s, "ターゲット市場規模", rx+Inches(0.2), Inches(1.65), Inches(6), Inches(0.3),
       fn=JP, sz=10, fc=GOLD)
    tb(s, "飲食+旅館・ホテル 約65万事業者 / うち事業承継課題 約25万",
       rx+Inches(0.2), Inches(1.98), Inches(6.1), Inches(0.62),
       fn=JP, sz=12, fc=WHITE, wrap=True)

    # KPI
    box(s, rx, Inches(3.0), Inches(6.43), Inches(2.3), fill=DARK3)
    tb(s, "年間KPI目標（テストマーケ期）", rx+Inches(0.2), Inches(3.12), Inches(6), Inches(0.3),
       fn=JP, sz=10, fc=GOLD)
    kpis = [
        ("月間LP訪問者数",     "200〜500 PV（6ヶ月以降）"),
        ("LINE登録数",         "月2〜5件（B2Bのため厳選）"),
        ("無料相談申込",       "月1〜2件"),
        ("成約（年間）",       "6〜10件（月1件ペース達成が目標）"),
    ]
    for j, (k, v) in enumerate(kpis):
        ky = Inches(3.48) + j * Inches(0.42)
        tb(s, k, rx+Inches(0.2), ky, Inches(3.0), Inches(0.38),
           fn=JP, sz=11, fc=RGBColor(0xAA,0xA2,0x95))
        tb(s, v, rx+Inches(3.4), ky, Inches(2.8), Inches(0.38),
           fn=JP, sz=11, fc=WHITE, bold=True)

    # LINE note
    box(s, rx, Inches(5.45), Inches(6.43), Inches(1.8), fill=RGBColor(0x18,0x2A,0x22))
    tb(s, "B2B集客の特性", rx+Inches(0.2), Inches(5.56), Inches(6), Inches(0.28),
       fn=JP, sz=9, fc=TEAL)
    tb(s, "一般消費者向けLINEと異なり、経営者へのリーチはセミナー・紹介が主体。\nLINEは温めた見込み客との関係維持ツールとして活用。月1件の成約が\nテストマーケティングの現実的な目標。",
       rx+Inches(0.2), Inches(5.86), Inches(6.1), Inches(1.2),
       fn=JP, sz=11, fc=RGBColor(0xAA,0xA2,0x95), wrap=True)


def s05_plans(prs):
    s = blank(prs)
    bg(s, CREAM)
    accent_bar(s)
    sec_label(s, "03 -- SERVICE PLANS", color=GOLD)
    heading(s, "サービスプラン比較", fc=DARK)
    gold_line(s, Inches(0.35), Inches(1.38))

    col_x = [Inches(0.25), Inches(3.45), Inches(6.7), Inches(9.95)]
    col_w = [Inches(3.1), Inches(3.15), Inches(3.15), Inches(3.13)]
    hbg   = [LGRAY, RGBColor(0xF0,0xED,0xE8), DARK2, DARK]
    htc   = [DARK, DARK, WHITE, WHITE]
    hpc   = [DARK, GOLD, GOLD_L, GOLD_L]
    names = ["", "ライト", "スタンダード", "プレミアム"]
    prices= ["", "150万円/年", "520万円/年", "910万円/年"]
    notes = ["", "うちLINE広告費 60万", "うちLINE広告費 60万\n+SNS制作費 100万", "うちLINE広告費 60万\n+SNS制作費 100万"]

    hh = Inches(1.2)
    for j in range(4):
        box(s, col_x[j], Inches(1.55), col_w[j], hh,
            fill=hbg[j], lc=BGRAY, lw=Pt(0.5))
        if j > 0:
            if j == 3:
                tb(s, "★ RECOMMEND", col_x[j]+Inches(0.1), Inches(1.6),
                   col_w[j]-Inches(0.15), Inches(0.25),
                   fn=JP, sz=8, fc=GOLD_L)
            tb(s, names[j], col_x[j]+Inches(0.12), Inches(1.82),
               col_w[j]-Inches(0.2), Inches(0.42), fn=JPS, sz=18, fc=htc[j])
            tb(s, prices[j], col_x[j]+Inches(0.12), Inches(2.22),
               col_w[j]-Inches(0.2), Inches(0.38), fn=JP, sz=14, fc=hpc[j], bold=True)
            tb(s, notes[j], col_x[j]+Inches(0.12), Inches(2.6),
               col_w[j]-Inches(0.2), Inches(0.4), fn=JP, sz=8,
               fc=RGBColor(0xAA,0xA0,0x90), wrap=True)

    rows = [
        ("支援期間",             ["2ヶ月",  "6ヶ月",   "12ヶ月"]),
        ("LP制作",               ["o",      "o",       "o"]),
        ("公式LINE構築",         ["基本",   "フル",    "フル"]),
        ("LINE広告 60万/年",     ["o",      "o",       "o"]),
        ("ブランディング戦略",   ["-",      "o",       "o"]),
        ("月次マーケ支援",       ["-",      "4ヶ月",   "12ヶ月"]),
        ("SNS動画制作・運用",    ["-",      "o 100万", "o 100万"]),
        ("AI活用・自動ブログ",   ["-",      "-",       "o"]),
        ("Brain 人件費（税抜）", ["90万",   "360万",   "750万"]),
    ]
    rbg = [WHITE, RGBColor(0xFA,0xF8,0xF4)]
    rh  = Inches(0.54)
    for ri, (label, vals) in enumerate(rows):
        ry = Inches(2.82) + ri * rh
        is_last = (ri == len(rows)-1)
        lbg = RGBColor(0x2A,0x27,0x22) if is_last else LGRAY
        box(s, col_x[0], ry, col_w[0], rh,
            fill=lbg, lc=BGRAY, lw=Pt(0.4))
        tb(s, label, col_x[0]+Inches(0.12), ry+Inches(0.12),
           col_w[0]-Inches(0.2), rh-Inches(0.12),
           fn=JP, sz=10.5, fc=GOLD_L if is_last else DARK, bold=is_last)

        vbg = [rbg[ri%2], DARK3, RGBColor(0x1C,0x1A,0x16)]
        if is_last: vbg = [RGBColor(0x1A,0x33,0x22)]*3
        for j, val in enumerate(vals):
            vr = RGBColor(0x99,0x90,0x88) if val == "-" else (GREEN if is_last else (GOLD_L if j==2 else (DARK if j==0 and not is_last else WHITE)))
            box(s, col_x[j+1], ry, col_w[j+1], rh,
                fill=vbg[j], lc=BGRAY, lw=Pt(0.4))
            disp = val.replace("o", "✓")
            tb(s, disp, col_x[j+1]+Inches(0.1), ry+Inches(0.12),
               col_w[j+1]-Inches(0.18), rh-Inches(0.12),
               fn=JP, sz=11, fc=vr, bold=is_last or (val != "-" and j==2),
               align=PP_ALIGN.CENTER)

    tb(s, "※SNS制作費・LINE広告費は実費扱いで別途計上。Brain人件費のみ補助金申請対象。",
       Inches(0.25), Inches(7.15), Inches(12.83), Inches(0.28),
       fn=JP, sz=8.5, fc=MGRAY)


def s06_lp(prs):
    s = blank(prs)
    bg(s, DARK)
    accent_bar(s)
    sec_label(s, "SERVICE DETAIL 01 -- LP / LANDING PAGE")
    heading(s, "LP（ランディングページ）制作")
    gold_line(s, Inches(0.35), Inches(1.38))

    # center message
    box(s, Inches(2.0), Inches(2.0), Inches(9.33), Inches(3.8),
        fill=DARK3, lc=GOLD, lw=Pt(0.75))
    tb(s, "LP は別途ご覧いただきます",
       Inches(2.0), Inches(2.6), Inches(9.33), Inches(0.7),
       fn=JPS, sz=28, fc=GOLD_L, align=PP_ALIGN.CENTER)
    tb(s, "現在制作中のLPをプロジェクターに投影してご説明します。",
       Inches(2.0), Inches(3.42), Inches(9.33), Inches(0.45),
       fn=JP, sz=14, fc=RGBColor(0xAA,0xA2,0x95), align=PP_ALIGN.CENTER)
    tb(s, "URL:  C:\\project\\daikonya\\index.html",
       Inches(2.0), Inches(4.0), Inches(9.33), Inches(0.38),
       fn=JP, sz=11, fc=RGBColor(0x66,0x60,0x58), align=PP_ALIGN.CENTER)

    # 含まれる内容
    items = ["和モダン × 高級感デザイン（Noto Serif JP）",
             "Hero / PAIN / MISSION / REASON / EXPERTISE / SERVICE / FAQ / CTA の全セクション",
             "スマートフォン対応（レスポンシブ対応）",
             "お問い合わせフォーム（無料相談申込）",
             "AI自動ブログ機能（プレミアムプランのみ、別途スライド参照）"]
    tb(s, "制作に含まれる内容", Inches(2.0), Inches(5.05), Inches(9.33), Inches(0.3),
       fn=JP, sz=9, fc=GOLD)
    for i, item in enumerate(items):
        ix = Inches(2.0) + (i % 2) * Inches(4.7)
        iy = Inches(5.38) + (i // 2) * Inches(0.45)
        box(s, ix, iy+Inches(0.14), Inches(0.07), Inches(0.07), fill=GOLD)
        tb(s, item, ix+Inches(0.2), iy+Inches(0.06), Inches(4.4), Inches(0.38),
           fn=JP, sz=11, fc=RGBColor(0xCC,0xC5,0xB8))


def s07_line(prs):
    s = blank(prs)
    bg(s, DARK)
    accent_bar(s, color=LINE_GREEN)
    sec_label(s, "SERVICE DETAIL 02 -- 公式LINE構築 + LINE広告")
    heading(s, "公式LINE構築 + LINE広告運用")
    gold_line(s, Inches(0.35), Inches(1.38))

    # LINE構築
    box(s, Inches(0.25), Inches(1.55), Inches(6.1), Inches(4.7),
        fill=DARK3, lc=LINE_GREEN, lw=Pt(1.0))
    tb(s, "公式LINE 構築", Inches(0.45), Inches(1.7), Inches(5.7), Inches(0.38),
       fn=JP, sz=13, fc=LINE_GREEN, bold=True)
    line_items = [
        "アカウント設定・プロフィール・アイコン",
        "リッチメニュー設計・作成",
        "自動挨拶メッセージ設定",
        "ステップ配信シナリオ設計（5〜8通）",
        "無料相談予約フロー構築",
        "セグメント配信設定（旅館 / 飲食 / 承継検討）",
    ]
    for i, item in enumerate(line_items):
        iy = Inches(2.18) + i * Inches(0.55)
        box(s, Inches(0.45), iy+Inches(0.16), Inches(0.08), Inches(0.08),
            fill=LINE_GREEN)
        tb(s, item, Inches(0.65), iy+Inches(0.08), Inches(5.5), Inches(0.42),
           fn=JP, sz=12, fc=RGBColor(0xCC,0xC5,0xB8))

    # LINE広告
    box(s, Inches(6.6), Inches(1.55), Inches(6.48), Inches(4.7),
        fill=RGBColor(0x18,0x26,0x18), lc=LINE_GREEN, lw=Pt(1.0))
    tb(s, "LINE広告 運用", Inches(6.8), Inches(1.7), Inches(6.0), Inches(0.38),
       fn=JP, sz=13, fc=LINE_GREEN, bold=True)
    tb(s, "年間 60万円", Inches(6.8), Inches(2.1), Inches(6.0), Inches(0.55),
       fn=JPS, sz=26, fc=GOLD_L)
    tb(s, "（広告費実費 / 全プラン共通）", Inches(6.8), Inches(2.65), Inches(6.0), Inches(0.3),
       fn=JP, sz=10, fc=RGBColor(0x88,0x80,0x70))
    ad_items = [
        "ターゲティング設定（旅館・飲食経営者層）",
        "クリエイティブ制作（バナー・テキスト）",
        "A/Bテスト実施・最適化",
        "月次広告レポート（CPL・CPAを追跡）",
        "LINE登録後のナーチャリング連携",
    ]
    for i, item in enumerate(ad_items):
        iy = Inches(3.1) + i * Inches(0.55)
        box(s, Inches(6.8), iy+Inches(0.16), Inches(0.08), Inches(0.08),
            fill=LINE_GREEN)
        tb(s, item, Inches(7.0), iy+Inches(0.08), Inches(6.0), Inches(0.42),
           fn=JP, sz=12, fc=RGBColor(0xCC,0xC5,0xB8))

    box(s, Inches(0.25), Inches(6.38), Inches(12.83), Inches(0.92),
        fill=RGBColor(0x18,0x26,0x18))
    tb(s, "全プランにLINE広告費 60万円/年を含みます。Brain管理費は各プランのBrain人件費に含みます。",
       Inches(0.45), Inches(6.52), Inches(12.4), Inches(0.62),
       fn=JP, sz=11, fc=RGBColor(0xAA,0xA2,0x95), wrap=True)


def s08_branding(prs):
    s = blank(prs)
    bg(s, DARK)
    accent_bar(s, color=GOLD)
    sec_label(s, "SERVICE DETAIL 03 -- ブランディング戦略支援")
    heading(s, "ブランディング戦略支援の内容")
    gold_line(s, Inches(0.35), Inches(1.38))

    items = [
        ("01", "ブランドコンセプト設計",
         "キャッチコピー・差別化軸・ブランドストーリーの言語化\n「なぜ小野氏に頼むべきか」を一言で言えるレベルに落とし込む"),
        ("02", "プロフィール・権威付けコンテンツ",
         "バイオグラフィ制作・実績の数字化・メディア露出戦略\nLP掲載用ケーススタディのフォーマット化"),
        ("03", "発信コンテンツカレンダー設計",
         "何を・いつ・どこで発信するかを3ヶ月単位で設計\nInstagram / note / セミナー を統合したテーマ管理"),
        ("04", "セミナー・登壇プランニング",
         "商工会・観光協会・金融機関との連携による登壇機会の創出\n年2〜4回の無料セミナー設計（認知→LINE登録への導線）"),
        ("05", "競合モニタリング",
         "類似コンサルタントの動向把握・差別化ポイントの継続確認\n四半期ごとにポジショニングを見直し"),
    ]
    cw = Inches(5.95)
    positions = [
        (Inches(0.25), Inches(1.55)),
        (Inches(6.5),  Inches(1.55)),
        (Inches(0.25), Inches(3.85)),
        (Inches(6.5),  Inches(3.85)),
        (Inches(0.25), Inches(6.05)),
    ]
    ch_list = [Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.1), Inches(1.2)]
    for i, ((num, title, desc), (bx, by)) in enumerate(zip(items, positions)):
        bh = ch_list[i]
        box(s, bx, by, cw if i < 4 else Inches(12.83), bh,
            fill=DARK3, lc=GOLD, lw=Pt(0.5))
        tb(s, num, bx+Inches(0.15), by+Inches(0.12), Inches(0.6), Inches(0.5),
           fn=JPS, sz=18, fc=RGBColor(0x55,0x50,0x40))
        tb(s, title, bx+Inches(0.75), by+Inches(0.15), (cw if i < 4 else Inches(12.83))-Inches(0.9), Inches(0.42),
           fn=JP, sz=13, fc=GOLD_L, bold=True)
        tb(s, desc, bx+Inches(0.75), by+Inches(0.55), (cw if i < 4 else Inches(12.83))-Inches(0.9), bh-Inches(0.65),
           fn=JP, sz=11, fc=RGBColor(0xAA,0xA2,0x95), wrap=True)


def s09_marketing(prs):
    s = blank(prs)
    bg(s, DARK)
    accent_bar(s, color=BLUE)
    sec_label(s, "SERVICE DETAIL 04 -- 月次マーケ支援")
    heading(s, "月次マーケティング支援")
    gold_line(s, Inches(0.35), Inches(1.38))

    # sub
    tb(s, "Google Analytics / Search Console の計測 + Web広告の運用成果を月次で分析・報告します。",
       Inches(0.35), Inches(1.5), Inches(12.6), Inches(0.38),
       fn=JP, sz=12, fc=RGBColor(0xAA,0xA2,0x95))

    tools = [
        ("Google Analytics 4",
         "LP訪問者数・流入経路・行動フロー・CVR を計測\nGoal設定（無料相談フォーム送信）の追跡",
         BLUE),
        ("Google Search Console",
         "検索クエリ・表示順位・CTR の可視化\nSEOキーワードの改善機会を月次で特定",
         RGBColor(0x34,0xA8,0x53)),
        ("Web広告（LINE / Google / Meta）",
         "広告インプレッション・CTR・CPL（1件当たりLINE登録コスト）\n月次予算配分の最適化提案",
         GOLD),
    ]
    for i, (tool, desc, color) in enumerate(tools):
        ty = Inches(2.05) + i * Inches(1.55)
        box(s, Inches(0.25), ty, Inches(12.83), Inches(1.42),
            fill=DARK3)
        box(s, Inches(0.25), ty, Inches(0.1), Inches(1.42), fill=color)
        tb(s, tool, Inches(0.5), ty+Inches(0.12), Inches(12.2), Inches(0.42),
           fn=JP, sz=13, fc=WHITE, bold=True)
        tb(s, desc, Inches(0.5), ty+Inches(0.56), Inches(12.2), Inches(0.72),
           fn=JP, sz=11, fc=RGBColor(0xAA,0xA2,0x95), wrap=True)

    # Report output
    box(s, Inches(0.25), Inches(6.78), Inches(12.83), Inches(0.56),
        fill=RGBColor(0x18,0x22,0x38))
    tb(s, "月次レポートのアウトプット：",
       Inches(0.45), Inches(6.88), Inches(3.0), Inches(0.32),
       fn=JP, sz=10, fc=BLUE)
    tb(s, "KPIダッシュボード（1枚）  /  改善提案（3点）  /  翌月アクションプラン（1枚）",
       Inches(3.6), Inches(6.88), Inches(9.2), Inches(0.32),
       fn=JP, sz=11, fc=WHITE)


def s10_sns(prs):
    s = blank(prs)
    bg(s, DARK)
    accent_bar(s, color=RGBColor(0xE1,0x30,0x6C))
    sec_label(s, "SERVICE DETAIL 05 -- SNS動画制作・運用")
    heading(s, "SNS動画制作・運用代行")
    gold_line(s, Inches(0.35), Inches(1.38))

    tb(s, "スタンダード・プレミアムプランに含みます　追加費用：100万円/年",
       Inches(0.35), Inches(1.5), Inches(12.6), Inches(0.38),
       fn=JP, sz=12, fc=GOLD_L, bold=True)

    # Monthly cycle
    cycle = [
        ("SHOOT", "現地撮影\n60分 / 月1回", RGBColor(0x8B,0x1A,0x4A)),
        ("EDIT",  "動画編集\n月8本×1分", RGBColor(0x6B,0x25,0x8C)),
        ("POST",  "SNS 投稿\nInstagram / YouTube", RGBColor(0x1A,0x5C,0x8C)),
        ("REPORT","月次分析\n再生数・フォロワー推移", RGBColor(0x1A,0x6B,0x3A)),
    ]
    cw2 = Inches(3.05)
    for i, (step, desc, color) in enumerate(cycle):
        cx = Inches(0.25) + i * (cw2 + Inches(0.12))
        box(s, cx, Inches(2.05), cw2, Inches(2.6), fill=color)
        tb(s, step, cx+Inches(0.15), Inches(2.18), cw2-Inches(0.25), Inches(0.45),
           fn=JP, sz=14, fc=WHITE, bold=True)
        tb(s, desc, cx+Inches(0.15), Inches(2.68), cw2-Inches(0.25), Inches(1.7),
           fn=JP, sz=16, fc=WHITE, wrap=True)
        if i < len(cycle)-1:
            tb(s, ">", cx+cw2, Inches(3.05), Inches(0.16), Inches(0.5),
               fn=JP, sz=14, fc=RGBColor(0x55,0x50,0x48), align=PP_ALIGN.CENTER)

    # Details
    details = [
        ("撮影", "月1回・60分の素材撮影（現地または リモート収録）\n撮影テーマは月次コンテンツカレンダーに基づき事前設計"),
        ("制作", "60分素材から月8本の1分Shorts動画を編集\nテロップ・BGM・サムネイル制作を含む"),
        ("投稿", "Instagram Reels / YouTube Shorts へ最適化して投稿代行\nハッシュタグ・説明文・ベストタイム投稿を管理"),
        ("効果測定", "再生回数・保存数・プロフィールへの流入数・フォロワー増加を月次報告"),
    ]
    for i, (lbl, desc) in enumerate(details):
        dy = Inches(4.85) + i * Inches(0.6)
        box(s, Inches(0.25), dy, Inches(1.4), Inches(0.52), fill=RGBColor(0x8B,0x1A,0x4A))
        tb(s, lbl, Inches(0.25), dy+Inches(0.1), Inches(1.4), Inches(0.35),
           fn=JP, sz=10, fc=WHITE, bold=True, align=PP_ALIGN.CENTER)
        box(s, Inches(1.7), dy, Inches(11.38), Inches(0.52), fill=DARK3)
        tb(s, desc, Inches(1.88), dy+Inches(0.07), Inches(11.0), Inches(0.42),
           fn=JP, sz=11, fc=RGBColor(0xCC,0xC5,0xB8), wrap=True)


def s11_ai(prs):
    s = blank(prs)
    bg(s, DARK)
    accent_bar(s, color=RGBColor(0x8B,0x5C,0xF6))
    sec_label(s, "SERVICE DETAIL 06 -- AI活用・自動ブログ")
    heading(s, "AI活用・自動ブログ実装")
    gold_line(s, Inches(0.35), Inches(1.38))
    tb(s, "プレミアムプランに含みます",
       Inches(0.35), Inches(1.5), Inches(8), Inches(0.35),
       fn=JP, sz=12, fc=RGBColor(0xAA,0x99,0xFF), bold=True)

    # Flow diagram
    flow = [
        ("GitHub Actions\n（定期実行）", RGBColor(0x24,0x29,0x3E)),
        ("AI（Claude API）\n記事自動生成",  RGBColor(0x2D,0x1F,0x4A)),
        ("LP ブログページ\n自動公開",        RGBColor(0x1A,0x2A,0x1A)),
    ]
    fw = Inches(3.8)
    for i, (label_t, color) in enumerate(flow):
        fx = Inches(0.25) + i * (fw + Inches(0.3))
        box(s, fx, Inches(2.0), fw, Inches(1.7), fill=color,
            lc=RGBColor(0x8B,0x5C,0xF6), lw=Pt(0.75))
        tb(s, label_t, fx+Inches(0.2), Inches(2.45), fw-Inches(0.35), Inches(0.9),
           fn=JP, sz=14, fc=WHITE, wrap=True, align=PP_ALIGN.CENTER)
        if i < len(flow)-1:
            tb(s, "->", fx+fw, Inches(2.6), Inches(0.34), Inches(0.5),
               fn=JP, sz=16, fc=RGBColor(0x8B,0x5C,0xF6), align=PP_ALIGN.CENTER)

    # Details
    details_ai = [
        ("自動ブログの仕組み",
         "GitHub Actions のcron設定（週1〜2回）でAI APIを呼び出し、旅館・飲食業に関する\n記事を自動生成。LPのブログセクションに自動投稿します。"),
        ("生成コンテンツのテーマ例",
         "観光DXの最新動向 / 補助金・助成金情報 / 事業承継の事例 / 旅館集客Tipsなど\n小野氏の専門性に沿ったトピックでSEO効果を狙います。"),
        ("ブログ機能の仕様",
         "LPに新セクション追加（Markdown対応）/ RSS配信対応 / 記事一覧・詳細ページ\nGitHub上で手動追記・編集も可能なシンプル構成。"),
        ("期待効果",
         "継続的なコンテンツ更新によるSEO改善 / 専門家としての権威付け強化\n更新の手間ゼロで「情報を発信し続ける専門家」としてのブランドを維持。"),
    ]
    for i, (dtitle, ddesc) in enumerate(details_ai):
        dy = Inches(3.88) + i * Inches(0.84)
        box(s, Inches(0.25), dy, Inches(12.83), Inches(0.76), fill=DARK3)
        box(s, Inches(0.25), dy, Inches(0.1), Inches(0.76),
            fill=RGBColor(0x8B,0x5C,0xF6))
        tb(s, dtitle, Inches(0.5), dy+Inches(0.08), Inches(3.5), Inches(0.3),
           fn=JP, sz=11, fc=RGBColor(0xAA,0x99,0xFF), bold=True)
        tb(s, ddesc, Inches(4.2), dy+Inches(0.06), Inches(8.7), Inches(0.64),
           fn=JP, sz=11, fc=RGBColor(0xAA,0xA2,0x95), wrap=True)


# ════════════════════════════════════════════════════════
#  MAIN
# ════════════════════════════════════════════════════════

def main():
    prs = new_prs()
    s01_cover(prs)
    s02_agenda(prs)
    s03_branding(prs)
    s04_funnel(prs)
    s05_plans(prs)
    s06_lp(prs)
    s07_line(prs)
    s08_branding(prs)
    s09_marketing(prs)
    s10_sns(prs)
    s11_ai(prs)
    prs.save(r"C:\project\daikonya\提案書_v2.pptx")
    print("Done: v2.pptx saved (11 slides)")

if __name__ == "__main__":
    main()
