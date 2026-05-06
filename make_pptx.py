"""
Generate two PPTX files for the Brain x TOJISM proposal meeting.
  1. 提案書.pptx         - Brain → 小野氏 への提案資料
  2. 補助金申請チェックシート.pptx - 申請に向けて決めること
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ──────────────────────────────────────────────
DARK   = RGBColor(0x1C, 0x1A, 0x16)
DARK2  = RGBColor(0x26, 0x23, 0x1E)
GOLD   = RGBColor(0xB8, 0x93, 0x5A)
GOLD_L = RGBColor(0xD4, 0xAF, 0x7A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CREAM  = RGBColor(0xF5, 0xF2, 0xED)
LGRAY  = RGBColor(0xEE, 0xEA, 0xE3)
MGRAY  = RGBColor(0x6B, 0x63, 0x59)
BGRAY  = RGBColor(0xDD, 0xD5, 0xC5)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x27, 0xAE, 0x60)

# ── Dimensions (16:9) ────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

JP  = "游ゴシック"
JPS = "游明朝"


# ── Helpers ─────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp

def txb(slide, text, l, t, w, h,
        font=JP, size=14, color=WHITE,
        bold=False, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    return box

def para(tf, text, font=JP, size=13, color=WHITE,
         bold=False, align=PP_ALIGN.LEFT, space_before=Pt(4)):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    return p

def label(slide, text, l, t, w=Inches(4), color=GOLD):
    txb(slide, text, l, t, w, Inches(0.3),
        font=JP, size=9, color=color)

def heading(slide, text, l, t, w, color=WHITE, size=28):
    txb(slide, text, l, t, w, Inches(1),
        font=JPS, size=size, color=color, bold=False)

def gold_bar(slide, l, t, w=Inches(0.35), h=Inches(0.03)):
    """Thin gold underline accent."""
    rect(slide, l, t, w, h, fill=GOLD)

def slide_bg(slide, color=DARK):
    rect(slide, 0, 0, SW, SH, fill=color)

def left_accent(slide, color=GOLD, w=Inches(0.12)):
    rect(slide, 0, 0, w, SH, fill=color)

def chip(slide, text, l, t, fill=DARK2, tc=GOLD_L, w=Inches(1.8), h=Inches(0.38)):
    r = rect(slide, l, t, w, h, fill=fill)
    txb(slide, text, l+Inches(0.12), t+Inches(0.06), w-Inches(0.24), h,
        font=JP, size=10, color=tc)
    return r

def row_box(slide, l, t, w, h, bg, texts, sizes, colors, bolds, aligns=None):
    """Utility to fill a rectangle with a list of text runs stacked."""
    rect(slide, l, t, w, h, fill=bg)
    if aligns is None:
        aligns = [PP_ALIGN.LEFT] * len(texts)
    box = slide.shapes.add_textbox(l+Inches(0.15), t+Inches(0.1),
                                   w-Inches(0.3), h-Inches(0.1))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(texts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(3)
        p.alignment = aligns[i]
        run = p.add_run()
        run.text = text
        run.font.name = JPS if sizes[i] >= 20 else JP
        run.font.size = Pt(sizes[i])
        run.font.color.rgb = colors[i]
        run.font.bold = bolds[i]


# ════════════════════════════════════════════════════════
#  FILE 1 : 提案書.pptx
# ════════════════════════════════════════════════════════

def make_proposal():
    prs = new_prs()

    # ── S1: 表紙 ────────────────────────────────────────
    s = blank(prs)
    slide_bg(s, DARK)
    # 左ゴールドバー
    rect(s, 0, 0, Inches(4.2), SH, fill=DARK2)
    rect(s, Inches(4.2), 0, Inches(0.06), SH, fill=GOLD)
    # タグライン
    txb(s, "Business Proposal  ·  2026.05.07",
        Inches(4.5), Inches(0.7), Inches(8), Inches(0.4),
        font=JP, size=10, color=RGBColor(0x88,0x80,0x70))
    # メインタイトル
    box = s.shapes.add_textbox(Inches(4.5), Inches(1.3), Inches(8.3), Inches(2.8))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "旅館・飲食業\nコンサルティング\nブランディング戦略"
    run.font.name = JPS; run.font.size = Pt(34)
    run.font.color.rgb = WHITE; run.font.bold = False
    # サブ
    txb(s, "Web戦略・集客設計・サポートプランのご提案",
        Inches(4.5), Inches(4.3), Inches(8), Inches(0.5),
        font=JP, size=13, color=RGBColor(0xAA,0xA0,0x90))
    # 左パネル情報
    txb(s, "提案先", Inches(0.4), Inches(1.8), Inches(3.2), Inches(0.35),
        font=JP, size=9, color=GOLD)
    txb(s, "小野 順平 氏\n株式会社TOJISM",
        Inches(0.4), Inches(2.15), Inches(3.2), Inches(0.9),
        font=JPS, size=16, color=WHITE)
    txb(s, "提案者", Inches(0.4), Inches(3.3), Inches(3.2), Inches(0.35),
        font=JP, size=9, color=GOLD)
    txb(s, "合同会社Brain\n鵜木 芳文",
        Inches(0.4), Inches(3.65), Inches(3.2), Inches(0.7),
        font=JP, size=13, color=WHITE)
    # ゴール
    rect(s, Inches(0.3), Inches(5.2), Inches(3.6), Inches(1.8),
         fill=RGBColor(0x28,0x26,0x22),
         line=GOLD, line_w=Pt(0.75))
    txb(s, "本日のゴール", Inches(0.5), Inches(5.35), Inches(3.2), Inches(0.35),
        font=JP, size=9, color=GOLD)
    txb(s, "ブランド方針に合意し\n補助金申請（5/22〆）\nに向けて動き出す",
        Inches(0.5), Inches(5.7), Inches(3.2), Inches(1.1),
        font=JP, size=12, color=WHITE)

    # ── S2: アジェンダ ──────────────────────────────────
    s = blank(prs)
    slide_bg(s, CREAM)
    left_accent(s, color=GOLD)
    label(s, "AGENDA", Inches(0.4), Inches(0.4), color=GOLD)
    heading(s, "本日のアジェンダ", Inches(0.4), Inches(0.75), Inches(6), color=DARK)
    gold_bar(s, Inches(0.4), Inches(1.35))

    items = [
        ("01", "ブランディング戦略",     "ポジショニング・ターゲット・差別化軸の確認",     "15 min", False),
        ("02", "Web戦略・集客ファネル",   "LP活用・公式LINE・SNS・マーケティング計画",        "20 min", False),
        ("03", "サポートプラン（松竹梅）", "3パターンの支援内容・期間・費用の比較",           "20 min", False),
        ("04", "補助金で実質負担ゼロへ",  "スキーム概要（詳細は別資料）",                  "10 min", False),
        ("05", "次のアクション",          "本日決めること・今週中に動くこと",               "15 min", True),
    ]
    for i, (num, title, sub, dur, urgent) in enumerate(items):
        y = Inches(1.65) + i * Inches(0.98)
        bg = RGBColor(0xFF,0xFA,0xF4) if not urgent else RGBColor(0xFF,0xF0,0xEE)
        rect(s, Inches(0.3), y, Inches(12.5), Inches(0.88), fill=bg,
             line=BGRAY, line_w=Pt(0.5))
        # 番号
        txb(s, num, Inches(0.5), y+Inches(0.1), Inches(0.6), Inches(0.7),
            font=JPS, size=20, color=GOLD)
        # タイトル
        txb(s, title, Inches(1.2), y+Inches(0.08), Inches(7.5), Inches(0.4),
            font=JP, size=14, color=DARK, bold=True)
        txb(s, sub, Inches(1.2), y+Inches(0.48), Inches(7.5), Inches(0.35),
            font=JP, size=11, color=MGRAY)
        # 時間バッジ
        bc = RED if urgent else RGBColor(0x99,0x90,0x88)
        rect(s, Inches(11.2), y+Inches(0.22), Inches(1.4), Inches(0.38),
             fill=bc)
        txb(s, ("★ 要確認 / " if urgent else "") + dur,
            Inches(11.22), y+Inches(0.3), Inches(1.36), Inches(0.3),
            font=JP, size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # ── S3: ポジショニング ───────────────────────────────
    s = blank(prs)
    slide_bg(s, CREAM)
    left_accent(s, color=GOLD)
    label(s, "01 — BRANDING STRATEGY", Inches(0.4), Inches(0.4), color=GOLD)
    heading(s, "ポジショニング：小野氏の立ち位置", Inches(0.4), Inches(0.75), Inches(9), color=DARK)
    gold_bar(s, Inches(0.4), Inches(1.35))

    cells = [
        ("一般コンサルタント",  "業種横断型・理論ベース\n現場経験は薄い",          CREAM,                  DARK,   MGRAY,  False),
        ("小野 順平 氏 ★",     "事業継承経験者\n旅館×飲食×AI/DX に対応できる\n数少ない専門人材", DARK, GOLD_L, WHITE,  True),
        ("旅館系コンサルタント","宿泊専門・飲食の知見薄\nDX活用は限定的",          RGBColor(0xF0,0xED,0xE8), DARK, MGRAY,  False),
        ("DX・IT系コンサルタント","技術特化・現場知識薄\n経営リアルが欠ける",      RGBColor(0xF0,0xED,0xE8), DARK, MGRAY,  False),
    ]
    positions = [
        (Inches(0.35),  Inches(1.55)),
        (Inches(6.85),  Inches(1.55)),
        (Inches(0.35),  Inches(4.35)),
        (Inches(6.85),  Inches(4.35)),
    ]
    cw, ch = Inches(6.2), Inches(2.6)
    for (bg, tc, sc, bold), (cx, cy) in zip(
            [(c[2],c[3],c[4],c[5]) for c in cells], positions):
        rect(s, cx, cy, cw, ch, fill=bg,
             line=GOLD if bold else BGRAY, line_w=Pt(1.5 if bold else 0.5))
    for (lbl, body, bg, tc, sc, bold), (cx, cy) in zip(cells, positions):
        txb(s, lbl, cx+Inches(0.2), cy+Inches(0.18), cw-Inches(0.4), Inches(0.4),
            font=JP, size=10, color=GOLD if bold else MGRAY)
        txb(s, body, cx+Inches(0.2), cy+Inches(0.55), cw-Inches(0.4), Inches(1.8),
            font=JP, size=13 if bold else 12, color=tc, bold=bold)

    # 差別化の3軸
    rect(s, Inches(0.35), Inches(7.0), Inches(12.62), Inches(0.38), fill=DARK2)
    txb(s, "差別化の3軸：  ① 事業継承の当事者経験（旅館2施設）　②  旅館×飲食 両業態への対応力　③ AI/ICT テクノロジーの実装支援",
        Inches(0.6), Inches(7.05), Inches(12.2), Inches(0.35),
        font=JP, size=10.5, color=GOLD_L)

    # ── S4: 集客ファネル ─────────────────────────────────
    s = blank(prs)
    slide_bg(s, DARK)
    left_accent(s, color=GOLD)
    label(s, "02 — WEB STRATEGY", Inches(0.4), Inches(0.4))
    heading(s, "集客ファネル設計", Inches(0.4), Inches(0.75), Inches(9))
    gold_bar(s, Inches(0.4), Inches(1.35))

    # ファネルボックス
    steps = [
        ("認知",          "Google検索 / Instagram / セミナー登壇 / 紹介・口コミ",         DARK2,  False),
        ("LP（ランディングページ）", "信頼構築・サービス説明・無料相談CTA",                RGBColor(0x28,0x26,0x22), True),
        ("公式LINE 登録", "無料相談予約・自動フォロー・情報配信",                         RGBColor(0x28,0x26,0x22), True),
        ("無料相談（60分）","課題ヒアリング → プラン提案",                               DARK2,  False),
        ("契約・支援スタート","スポット or 継続プラン",                                   DARK2,  False),
    ]
    sx = Inches(0.35); bw = Inches(5.8); bh = Inches(0.82)
    for i, (title, sub, bg, highlight) in enumerate(steps):
        sy = Inches(1.6) + i * (bh + Inches(0.18))
        lc = GOLD if highlight else RGBColor(0x44,0x40,0x38)
        rect(s, sx, sy, bw, bh, fill=bg,
             line=lc, line_w=Pt(1.5 if highlight else 0.5))
        tc = GOLD_L if highlight else RGBColor(0xCC,0xC5,0xB8)
        txb(s, title, sx+Inches(0.18), sy+Inches(0.1), bw-Inches(0.3), Inches(0.4),
            font=JP, size=14, color=tc, bold=highlight)
        txb(s, sub, sx+Inches(0.18), sy+Inches(0.46), bw-Inches(0.3), Inches(0.3),
            font=JP, size=10, color=RGBColor(0x88,0x80,0x70))
        if i < len(steps)-1:
            txb(s, "↓", sx+Inches(2.6), sy+bh, Inches(0.6), Inches(0.25),
                font=JP, size=11, color=GOLD, align=PP_ALIGN.CENTER)

    # 右側のポイント
    points = [
        ("初期集客の優先順位",
         "① 紹介・人脈（最速）\n② セミナー登壇\n③ Instagram 経営Tips発信\n④ Google SEO（中長期）"),
        ("LINEを中間に置く理由",
         "旅館・飲食オーナーはメール開封率が低い。\nLINEなら既読率が高く、無料コンテンツで\n継続的な接点を維持できる。"),
        ("年間KPI目安",
         "LINE登録 月10〜20人\n無料相談 月2〜4件\n成約率 30〜40%（月1件ペース）"),
    ]
    rx = Inches(6.5)
    for i, (pt, desc) in enumerate(points):
        ry = Inches(1.6) + i * Inches(1.9)
        rect(s, rx, ry, Inches(6.5), Inches(1.7),
             fill=RGBColor(0x22,0x20,0x1C))
        txb(s, pt, rx+Inches(0.2), ry+Inches(0.12), Inches(6.1), Inches(0.38),
            font=JP, size=11, color=GOLD, bold=True)
        txb(s, desc, rx+Inches(0.2), ry+Inches(0.48), Inches(6.1), Inches(1.1),
            font=JP, size=11, color=RGBColor(0xAA,0xA2,0x95))

    # ── S5: 3プラン比較 ──────────────────────────────────
    s = blank(prs)
    slide_bg(s, CREAM)
    left_accent(s, color=GOLD)
    label(s, "03 — SERVICE PLANS", Inches(0.4), Inches(0.4), color=GOLD)
    heading(s, "サポートプラン（松・竹・梅）比較", Inches(0.4), Inches(0.75), Inches(9), color=DARK)
    gold_bar(s, Inches(0.4), Inches(1.35))

    # ヘッダー
    cols = [Inches(0.3), Inches(3.5), Inches(6.75), Inches(10.0)]
    cw_plan = Inches(3.1)
    hbgs = [CREAM, RGBColor(0xF0,0xED,0xE8), DARK2, DARK]
    hnames = ["", "梅  スターター", "竹  スタンダード", "松  フルサポート ★"]
    hprices= ["", "90万円（税抜）", "360万円（税抜）", "750万円（税抜）"]
    hpd    = ["", "/ 一括", "/ 分割可", "/ 補助金フル活用"]
    htc    = [DARK, DARK, WHITE, WHITE]
    hpc    = [DARK, GOLD, GOLD_L, GOLD_L]

    hrow_h = Inches(1.1)
    for i in range(4):
        rect(s, cols[i], Inches(1.55), cw_plan if i>0 else Inches(3.1),
             hrow_h, fill=hbgs[i],
             line=BGRAY, line_w=Pt(0.5))
        if i > 0:
            txb(s, hnames[i], cols[i]+Inches(0.12), Inches(1.62), cw_plan-Inches(0.2), Inches(0.4),
                font=JPS, size=15, color=htc[i], bold=False)
            txb(s, hprices[i], cols[i]+Inches(0.12), Inches(2.0), cw_plan-Inches(0.2), Inches(0.42),
                font=JPS, size=18, color=hpc[i], bold=False)
            txb(s, hpd[i], cols[i]+Inches(0.12), Inches(2.45), cw_plan-Inches(0.2), Inches(0.2),
                font=JP, size=9, color=RGBColor(0xAA,0xA0,0x90) if i>1 else MGRAY)

    # 行データ
    rows = [
        ("支援期間",            ["2ヶ月",   "6ヶ月",   "12ヶ月"],    [False, False, True]),
        ("LP制作（本格版）",    ["✓",       "✓",       "✓"],          [False, False, False]),
        ("公式LINE構築",        ["基本設定", "✓",       "✓（フル）"],  [False, False, True]),
        ("ブランディング戦略",  ["簡易版",  "✓",       "✓（フル）"],  [False, False, True]),
        ("月次マーケ支援",      ["—",       "4ヶ月",   "12ヶ月"],     [False, False, True]),
        ("SNS / 広告運用",      ["—",       "—",       "✓"],          [False, False, True]),
        ("補助金申請額",        ["100万円", "400万円", "800万円"],    [False, False, True]),
        ("小野氏の実質負担",    ["ほぼ0円", "ほぼ0円", "ほぼ0円"],   [True, True, True]),
    ]
    rbg_alt  = [WHITE, RGBColor(0xFA,0xF8,0xF4)]
    row_h = Inches(0.52)
    for ri, (lbl, vals, bolds) in enumerate(rows):
        ry = Inches(2.75) + ri * row_h
        is_last = (ri == len(rows)-1)
        lbg = RGBColor(0x2A,0x27,0x22) if is_last else LGRAY
        rect(s, cols[0], ry, Inches(3.1), row_h, fill=lbg,
             line=BGRAY, line_w=Pt(0.5))
        lc = GOLD_L if is_last else DARK
        txb(s, lbl, cols[0]+Inches(0.12), ry+Inches(0.12), Inches(2.8), row_h-Inches(0.1),
            font=JP, size=11, color=lc, bold=is_last)

        val_bgs = [rbg_alt[ri%2], RGBColor(0x24,0x22,0x1E), RGBColor(0x1C,0x1A,0x16)]
        val_tcs = [MGRAY if vals[j]=="—" else DARK if j==0 else WHITE for j in range(3)]
        for j in range(3):
            vbg = RGBColor(0x22,0x44,0x2E) if is_last else val_bgs[j]
            rect(s, cols[j+1], ry, cw_plan, row_h, fill=vbg,
                 line=BGRAY, line_w=Pt(0.5))
            vc = GREEN if is_last else (GOLD_L if bolds[j] else val_tcs[j])
            txb(s, vals[j], cols[j+1]+Inches(0.12), ry+Inches(0.12),
                cw_plan-Inches(0.2), row_h-Inches(0.1),
                font=JP, size=12 if is_last else 11, color=vc,
                bold=bolds[j] or is_last, align=PP_ALIGN.CENTER)

    txb(s, "※補助金フル採択の場合。消費税相当分は課税区分により別途考慮が必要です。",
        Inches(0.35), Inches(7.12), Inches(12.6), Inches(0.3),
        font=JP, size=9, color=MGRAY)

    # ── S6: 補助金で実質ゼロ ────────────────────────────
    s = blank(prs)
    slide_bg(s, DARK)
    left_accent(s, color=GOLD)
    label(s, "04 — SUBSIDY", Inches(0.4), Inches(0.4))
    heading(s, "補助金活用で、実質負担ほぼゼロへ", Inches(0.4), Inches(0.75), Inches(10))
    gold_bar(s, Inches(0.4), Inches(1.35))

    # 大きな数字
    rect(s, Inches(0.35), Inches(1.6), Inches(5.5), Inches(4.8),
         fill=RGBColor(0x22,0x20,0x1C))
    txb(s, "小野氏の実質負担", Inches(0.6), Inches(1.85), Inches(5), Inches(0.4),
        font=JP, size=11, color=GOLD)
    txb(s, "ほぼ", Inches(0.6), Inches(2.3), Inches(2), Inches(0.8),
        font=JP, size=20, color=WHITE)
    txb(s, "0円", Inches(0.6), Inches(2.95), Inches(5), Inches(1.4),
        font=JPS, size=64, color=GREEN, bold=False)
    txb(s, "（観光DX補助金・定額補助 上限800万円を活用）",
        Inches(0.6), Inches(4.5), Inches(5), Inches(0.5),
        font=JP, size=10, color=RGBColor(0x88,0x80,0x70))

    # 右側：条件・注意
    rx = Inches(6.3)
    conds = [
        ("上限800万円の定額補助",      "かかった費用がそのまま補助されます。\n1/2補助や2/3補助ではありません。"),
        ("専門人材の人件費・交通費対象", "LP制作・LINE・Web戦略支援は\n「観光DX」の文脈で申請します。"),
        ("後払いに注意",               "補助金入金は2027年2〜3月頃。\n先払いが必要なため資金計画が重要です。\n（つなぎ融資の活用も可）"),
    ]
    for i, (ct, cd) in enumerate(conds):
        cy = Inches(1.55) + i * Inches(1.85)
        rect(s, rx, cy, Inches(6.7), Inches(1.7),
             fill=RGBColor(0x26,0x23,0x1E))
        if i == 2:
            rect(s, rx, cy, Inches(0.08), Inches(1.7), fill=RED)
        else:
            rect(s, rx, cy, Inches(0.08), Inches(1.7), fill=GOLD)
        txb(s, ct, rx+Inches(0.22), cy+Inches(0.14), Inches(6.2), Inches(0.38),
            font=JP, size=12, color=GOLD_L, bold=True)
        txb(s, cd, rx+Inches(0.22), cy+Inches(0.52), Inches(6.2), Inches(1.05),
            font=JP, size=11, color=RGBColor(0xAA,0xA2,0x95))

    txb(s, "詳細は「補助金申請チェックシート」をご参照ください。",
        Inches(6.3), Inches(7.12), Inches(6.7), Inches(0.3),
        font=JP, size=9.5, color=RGBColor(0x66,0x60,0x58))

    # ── S7: 次のアクション ──────────────────────────────
    s = blank(prs)
    slide_bg(s, DARK)
    left_accent(s, color=GOLD)
    label(s, "05 — NEXT ACTIONS", Inches(0.4), Inches(0.4))
    heading(s, "本日決めること・今週中に動くこと", Inches(0.4), Inches(0.75), Inches(10))
    gold_bar(s, Inches(0.4), Inches(1.35))

    actions = [
        ("小野氏", "プランの選択（松・竹・梅）",  "補助金申請額・支援期間・内容の最終確認",  "本日中",          True),
        ("小野氏", "補助金 参加申込",            "kanko-dx-hojo.go.jp から申込",            "5月22日（木）〆", True),
        ("Brain", "計画申請書の作成サポート",    "DX文脈での事業計画書を共同で起草",        "5月29日（木）〆", True),
        ("小野氏", "コンテンツ素材の提供",       "MISSIONテキスト・実績・経歴（LP用）",     "5月中",           False),
        ("Brain", "業務委託契約書の準備",        "補助金採択後の正式契約に向けた事前合意",   "6月上旬",         False),
    ]
    for i, (who, title, sub, deadline, urgent) in enumerate(actions):
        ay = Inches(1.6) + i * Inches(1.06)
        rect(s, Inches(0.35), ay, Inches(12.62), Inches(0.96),
             fill=RGBColor(0x22,0x20,0x1C))
        # who badge
        wc = GOLD if who == "小野氏" else GREEN
        rect(s, Inches(0.35), ay, Inches(1.2), Inches(0.96), fill=wc)
        txb(s, who, Inches(0.37), ay+Inches(0.3), Inches(1.16), Inches(0.38),
            font=JP, size=11, color=DARK, bold=True, align=PP_ALIGN.CENTER)
        txb(s, title, Inches(1.7), ay+Inches(0.08), Inches(7.5), Inches(0.42),
            font=JP, size=13, color=WHITE, bold=False)
        txb(s, sub,  Inches(1.7), ay+Inches(0.52), Inches(7.5), Inches(0.38),
            font=JP, size=10, color=RGBColor(0x88,0x80,0x70))
        dc = RED if urgent else RGBColor(0x66,0x60,0x58)
        txb(s, deadline, Inches(10.5), ay+Inches(0.28), Inches(2.3), Inches(0.42),
            font=JP, size=11, color=dc, bold=urgent, align=PP_ALIGN.RIGHT)

    # ハイライト
    rect(s, Inches(0.35), Inches(7.0), Inches(12.62), Inches(0.42),
         fill=RGBColor(0x3A,0x18,0x14))
    txb(s, "🔑  参加申込（5/22〆）を今週中に完了することが「実質負担ゼロ」実現の唯一の条件です。",
        Inches(0.55), Inches(7.06), Inches(12.2), Inches(0.32),
        font=JP, size=11, color=RGBColor(0xFF,0xAA,0xAA))

    prs.save(r"C:\project\daikonya\提案書.pptx")
    print("OK: teian-sho.pptx saved")


# ════════════════════════════════════════════════════════
#  FILE 2 : 補助金申請チェックシート.pptx
# ════════════════════════════════════════════════════════

def make_subsidy():
    prs = new_prs()

    # ── S1: 表紙 ────────────────────────────────────────
    s = blank(prs)
    slide_bg(s, DARK2)
    rect(s, 0, 0, Inches(0.5), SH, fill=RED)
    txb(s, "補助金申請チェックシート",
        Inches(0.8), Inches(1.8), Inches(10), Inches(1.2),
        font=JPS, size=36, color=WHITE, bold=False)
    txb(s, "観光DX推進事業 専門人材による伴走支援",
        Inches(0.8), Inches(3.1), Inches(10), Inches(0.5),
        font=JP, size=16, color=RGBColor(0xCC,0xC5,0xB8))
    txb(s, "申請に向けて決めること・確認事項をまとめた作業資料です",
        Inches(0.8), Inches(3.75), Inches(10), Inches(0.4),
        font=JP, size=12, color=RGBColor(0x88,0x80,0x70))
    rect(s, Inches(0.8), Inches(5.2), Inches(7), Inches(1.6),
         fill=RGBColor(0x3A,0x18,0x14))
    txb(s, "⚠  締切厳守", Inches(1.0), Inches(5.35), Inches(3), Inches(0.38),
        font=JP, size=11, color=RED, bold=True)
    txb(s, "参加申込：2026年5月22日（金）17:00\n計画申請：2026年5月29日（金）17:00",
        Inches(1.0), Inches(5.72), Inches(6.5), Inches(0.88),
        font=JP, size=14, color=RGBColor(0xFF,0xAA,0xAA))
    txb(s, "作成：合同会社Brain  /  2026年5月7日",
        Inches(0.8), Inches(7.1), Inches(8), Inches(0.3),
        font=JP, size=9, color=RGBColor(0x66,0x60,0x58))

    # ── S2: 補助金概要 ───────────────────────────────────
    s = blank(prs)
    slide_bg(s, DARK)
    left_accent(s, color=RED)
    label(s, "CHECK 01 — 補助金概要の確認", Inches(0.4), Inches(0.4))
    heading(s, "観光DX補助金とは何か", Inches(0.4), Inches(0.75), Inches(9))
    gold_bar(s, Inches(0.4), Inches(1.35), w=Inches(0.32), h=Inches(0.03))

    # 4つの数字
    stats = [
        ("800万円", "補助上限額（定額）"),
        ("11,800円", "最高時間単価\n（税込 / 時間）"),
        ("5月22日", "参加申込 締切"),
        ("5月29日", "計画申請 締切"),
    ]
    for i, (num, lbl) in enumerate(stats):
        sx = Inches(0.35) + i * Inches(3.24)
        bc = RED if i >= 2 else RGBColor(0x26,0x23,0x1E)
        rect(s, sx, Inches(1.55), Inches(3.1), Inches(1.55), fill=bc)
        txb(s, num, sx+Inches(0.15), Inches(1.68), Inches(2.8), Inches(0.85),
            font=JPS, size=26, color=WHITE if i >= 2 else GOLD_L)
        txb(s, lbl, sx+Inches(0.15), Inches(2.52), Inches(2.8), Inches(0.48),
            font=JP, size=10, color=RGBColor(0xFF,0xBB,0xBB) if i >= 2 else RGBColor(0x88,0x80,0x70))

    # 補助対象
    rect(s, Inches(0.35), Inches(3.3), Inches(12.62), Inches(0.38),
         fill=RGBColor(0x26,0x23,0x1E))
    txb(s, "補助対象となる活動（観光DXの文脈で申請する）",
        Inches(0.5), Inches(3.36), Inches(12), Inches(0.28),
        font=JP, size=10, color=GOLD)

    targets = [
        "観光DXに関する計画の策定（LPをDX集客戦略計画として位置付け）",
        "旅行者の利便性向上に資するデジタルツールの導入（公式LINE・予約システム連携）",
        "デジタルツール導入後の活用支援（月次Webマーケ・データ分析・SNS運用）",
    ]
    for i, t in enumerate(targets):
        ry = Inches(3.82) + i * Inches(0.72)
        rect(s, Inches(0.35), ry, Inches(0.42), Inches(0.52), fill=GOLD)
        txb(s, str(i+1), Inches(0.35), ry+Inches(0.08), Inches(0.42), Inches(0.38),
            font=JP, size=13, color=DARK, bold=True, align=PP_ALIGN.CENTER)
        rect(s, Inches(0.82), ry, Inches(12.1), Inches(0.52),
             fill=RGBColor(0x22,0x20,0x1C))
        txb(s, t, Inches(1.0), ry+Inches(0.1), Inches(11.8), Inches(0.36),
            font=JP, size=12, color=RGBColor(0xCC,0xC5,0xB8))

    rect(s, Inches(0.35), Inches(6.2), Inches(12.62), Inches(1.1),
         fill=RGBColor(0x28,0x26,0x22))
    txb(s, "補助金の特徴", Inches(0.55), Inches(6.3), Inches(3), Inches(0.3),
        font=JP, size=9, color=GOLD)
    txb(s, "定額補助（1/2や2/3ではなく全額補填）　／　上限800万円　／　後払い方式（先に支払い→後から入金）　／　時間単価の証憑が必要",
        Inches(0.55), Inches(6.62), Inches(12.2), Inches(0.55),
        font=JP, size=11, color=RGBColor(0xAA,0xA2,0x95))

    # ── S3: 申請者の確認 ────────────────────────────────
    s = blank(prs)
    slide_bg(s, DARK)
    left_accent(s, color=RED)
    label(s, "CHECK 02 — 申請者の確認  ★ 要決定", Inches(0.4), Inches(0.4))
    heading(s, "誰が申請者（補助対象事業者）になるか？", Inches(0.4), Inches(0.75), Inches(11))
    gold_bar(s, Inches(0.4), Inches(1.35), w=Inches(0.32), h=Inches(0.03))

    txb(s, "補助対象事業者の条件：地方公共団体 / DMO / 観光協会 / 宿泊事業者等",
        Inches(0.35), Inches(1.55), Inches(12.62), Inches(0.38),
        font=JP, size=11, color=RGBColor(0xAA,0xA2,0x95))

    options = [
        ("案 A", "丸正旅館（または旅館部門）",
         "宿泊事業者として明確に適格。\n本店所在地が丸正旅館内のため関係性も説明しやすい。",
         "旅館の代表者（オーナー）が申請者となるため\n小野氏の関与方法を整理する必要あり。", GOLD),
        ("案 B", "株式会社TOJISM",
         "小野氏自身の法人。自由に動きやすい。",
         "宿泊事業者等として認められるか確認が必要。\nコンサル会社としての申請は審査が厳しい場合あり。",
         RGBColor(0x99,0x90,0x80)),
    ]
    for i, (plan, name, merit, risk, ac) in enumerate(options):
        ox = Inches(0.35) + i * Inches(6.45)
        rect(s, ox, Inches(2.1), Inches(6.2), Inches(4.6),
             fill=RGBColor(0x22,0x20,0x1C),
             line=ac, line_w=Pt(1.5 if i==0 else 0.5))
        txb(s, plan, ox+Inches(0.2), Inches(2.22), Inches(5.8), Inches(0.35),
            font=JP, size=10, color=ac)
        txb(s, name, ox+Inches(0.2), Inches(2.55), Inches(5.8), Inches(0.55),
            font=JP, size=15, color=WHITE, bold=True)
        txb(s, "✓ メリット", ox+Inches(0.2), Inches(3.18), Inches(5.8), Inches(0.3),
            font=JP, size=10, color=GREEN)
        txb(s, merit, ox+Inches(0.2), Inches(3.5), Inches(5.8), Inches(0.75),
            font=JP, size=11, color=RGBColor(0xAA,0xA2,0x95))
        txb(s, "⚠ 確認事項", ox+Inches(0.2), Inches(4.35), Inches(5.8), Inches(0.3),
            font=JP, size=10, color=RED)
        txb(s, risk, ox+Inches(0.2), Inches(4.67), Inches(5.8), Inches(0.75),
            font=JP, size=11, color=RGBColor(0xAA,0xA2,0x95))
        if i == 0:
            rect(s, ox, Inches(2.1), Inches(1.2), Inches(0.45), fill=GOLD)
            txb(s, "← 推奨", ox+Inches(0.06), Inches(2.22), Inches(1.08), Inches(0.3),
                font=JP, size=10, color=DARK, bold=True)

    rect(s, Inches(0.35), Inches(6.85), Inches(12.62), Inches(0.5),
         fill=RGBColor(0x3A,0x18,0x14))
    txb(s, "□  本日確認：どちらの法人を申請者とするか決定してください。",
        Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.32),
        font=JP, size=11, color=RGBColor(0xFF,0xCC,0xCC))

    # ── S4: プランと費用確認 ────────────────────────────
    s = blank(prs)
    slide_bg(s, DARK)
    left_accent(s, color=RED)
    label(s, "CHECK 03 — プラン & 費用の確認  ★ 要決定", Inches(0.4), Inches(0.4))
    heading(s, "申請プランと費用シミュレーション", Inches(0.4), Inches(0.75), Inches(11))
    gold_bar(s, Inches(0.4), Inches(1.35), w=Inches(0.32), h=Inches(0.03))

    # テーブル
    thead = ["項目", "梅（スターター）", "竹（スタンダード）", "松（フルサポート） ★"]
    trows = [
        ["期間",         "2ヶ月",           "6ヶ月",             "12ヶ月"],
        ["Brain費用（税抜）","90万円",        "360万円",          "750万円"],
        ["時間数（目安）","約80時間",        "約330時間",         "約670時間"],
        ["補助金申請額", "100万円",          "400万円",           "800万円（上限）"],
        ["小野氏負担",   "ほぼ0円",          "ほぼ0円",           "ほぼ0円"],
    ]
    col_x = [Inches(0.35), Inches(3.15), Inches(6.35), Inches(9.55)]
    col_w = [Inches(2.7), Inches(3.1), Inches(3.1), Inches(3.73)]
    # ヘッダー
    hbg2 = [LGRAY, RGBColor(0xF0,0xED,0xE8), RGBColor(0x2A,0x27,0x22), RGBColor(0x1C,0x1A,0x16)]
    htc2 = [DARK, DARK, WHITE, WHITE]
    hlines = [BGRAY, BGRAY, RGBColor(0x55,0x50,0x48), GOLD]
    for j in range(4):
        rect(s, col_x[j], Inches(1.55), col_w[j], Inches(0.55),
             fill=hbg2[j], line=hlines[j], line_w=Pt(0.75 if j==3 else 0.5))
        txb(s, thead[j], col_x[j]+Inches(0.12), Inches(1.63), col_w[j]-Inches(0.2), Inches(0.38),
            font=JP, size=11, color=htc2[j], bold=True,
            align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)
    # 行
    for ri, row in enumerate(trows):
        ry = Inches(2.16) + ri * Inches(0.9)
        is_last = ri == len(trows)-1
        row_alt = [WHITE, RGBColor(0xFA,0xF8,0xF4)]
        rbg = [LGRAY, row_alt[ri%2], RGBColor(0x22,0x20,0x1C), RGBColor(0x1C,0x1A,0x16)]
        if is_last: rbg = [RGBColor(0x22,0x44,0x2E)]*4
        for j in range(4):
            rect(s, col_x[j], ry, col_w[j], Inches(0.85),
                 fill=rbg[j], line=BGRAY, line_w=Pt(0.5))
            vc = DARK if (j==0 and not is_last) else (GREEN if is_last else (GOLD_L if j==3 else WHITE))
            txb(s, row[j], col_x[j]+Inches(0.12), ry+Inches(0.2),
                col_w[j]-Inches(0.2), Inches(0.5),
                font=JP, size=12, color=vc, bold=is_last or j==3,
                align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)

    rect(s, Inches(0.35), Inches(6.82), Inches(12.62), Inches(0.5),
         fill=RGBColor(0x3A,0x18,0x14))
    txb(s, "□  本日確認：松・竹・梅のどのプランにするかを決定してください。",
        Inches(0.55), Inches(6.92), Inches(12.2), Inches(0.32),
        font=JP, size=11, color=RGBColor(0xFF,0xCC,0xCC))

    # ── S5: キャッシュフロー確認 ────────────────────────
    s = blank(prs)
    slide_bg(s, DARK)
    left_accent(s, color=RED)
    label(s, "CHECK 04 — キャッシュフロー確認  ★ 要確認", Inches(0.4), Inches(0.4))
    heading(s, "先払いが必要 ── 資金手当の確認", Inches(0.4), Inches(0.75), Inches(11))
    gold_bar(s, Inches(0.4), Inches(1.35), w=Inches(0.32), h=Inches(0.03))

    timeline_items = [
        ("5月", "参加申込\n計画申請",       RED,   "申請完了"),
        ("6月", "採択審査\n交付申請",       MGRAY, "待機期間"),
        ("7月", "交付決定\n発注開始",       GOLD,  "ここから\nBrainへ支払い"),
        ("8〜12月", "事業実施\n月次支払い",  RGBColor(0x27,0x6B,0x3F), "立替期間\n（最大750万円）"),
        ("翌1月", "完了報告",              MGRAY, "書類提出"),
        ("翌3月", "補助金\n入金",          GREEN, "資金回収"),
    ]
    tw = Inches(2.0)
    for i, (month, label_t, color, note) in enumerate(timeline_items):
        tx = Inches(0.35) + i * Inches(2.18)
        rect(s, tx, Inches(1.65), tw, Inches(0.7), fill=color)
        txb(s, month, tx, Inches(1.72), tw, Inches(0.56),
            font=JP, size=13, color=DARK if color in [GOLD, GREEN, RGBColor(0x27,0x6B,0x3F)] else WHITE,
            bold=True, align=PP_ALIGN.CENTER)
        rect(s, tx, Inches(2.42), tw, Inches(1.4),
             fill=RGBColor(0x22,0x20,0x1C))
        txb(s, label_t, tx+Inches(0.1), Inches(2.52), tw-Inches(0.2), Inches(0.8),
            font=JP, size=10, color=RGBColor(0xCC,0xC5,0xB8), align=PP_ALIGN.CENTER)
        txb(s, note, tx+Inches(0.05), Inches(3.9), tw-Inches(0.1), Inches(0.7),
            font=JP, size=9, color=RGBColor(0x88,0x80,0x70), align=PP_ALIGN.CENTER)
        if i < len(timeline_items)-1:
            txb(s, "→", Inches(0.35) + i * Inches(2.18) + tw, Inches(1.88),
                Inches(0.18), Inches(0.38),
                font=JP, size=14, color=RGBColor(0x55,0x50,0x48), align=PP_ALIGN.CENTER)

    rect(s, Inches(0.35), Inches(4.8), Inches(12.62), Inches(0.5),
         fill=RGBColor(0x28,0x26,0x22))
    txb(s, "立替が必要な最大金額",
        Inches(0.55), Inches(4.88), Inches(4), Inches(0.3),
        font=JP, size=10, color=GOLD)

    plans_cf = [
        ("梅プランの場合", "最大 約100万円 の立替"),
        ("竹プランの場合", "最大 約400万円 の立替"),
        ("松プランの場合", "最大 約800万円 の立替"),
    ]
    for i, (pl, cf) in enumerate(plans_cf):
        px = Inches(0.35) + i * Inches(4.25)
        rect(s, px, Inches(5.38), Inches(4.1), Inches(1.28),
             fill=RGBColor(0x22,0x20,0x1C))
        txb(s, pl, px+Inches(0.15), Inches(5.47), Inches(3.8), Inches(0.35),
            font=JP, size=10, color=MGRAY)
        txb(s, cf, px+Inches(0.15), Inches(5.8), Inches(3.8), Inches(0.65),
            font=JP, size=13, color=WHITE, bold=True)

    rect(s, Inches(0.35), Inches(6.78), Inches(12.62), Inches(0.55),
         fill=RGBColor(0x26,0x23,0x1E))
    txb(s, "対策：日本政策金融公庫の「つなぎ融資」/ Brainとの月次分割払い契約 / 自己資金での対応",
        Inches(0.55), Inches(6.89), Inches(12.2), Inches(0.35),
        font=JP, size=10.5, color=RGBColor(0xAA,0xA2,0x95))

    # ── S6: 必要書類チェックリスト ──────────────────────
    s = blank(prs)
    slide_bg(s, DARK)
    left_accent(s, color=RED)
    label(s, "CHECK 05 — 必要書類", Inches(0.4), Inches(0.4))
    heading(s, "申請に必要な書類チェックリスト", Inches(0.4), Inches(0.75), Inches(11))
    gold_bar(s, Inches(0.4), Inches(1.35), w=Inches(0.32), h=Inches(0.03))

    docs = [
        ("必須", "【様式1-1】計画申請書",           "Brain（専門人材）ごとに作成"),
        ("必須", "【様式1-2】スケジュール・実施体制", "支援内容と工程表"),
        ("必須", "【様式1-3】専門人材の同意書",      "Brain代表が署名"),
        ("必須", "【様式2】補助対象経費内訳",        "費用の内訳を記載"),
        ("必須", "【様式3】補助対象経費算定根拠",    "交通費等の証憑も添付"),
        ("必須", "見積書",                          "Brain作成。時間単価・作業時間を明記"),
        ("必須", "時間単価の算定根拠資料",           "Brainの過去類似業務の時間単価証憑"),
        ("必須", "【様式4】チェックリスト",          "申請内容に不備がないか確認"),
        ("条件付", "申請団体の存立証明書類",         "観光協会の場合は必須（登記事項証明書等）"),
    ]
    for i, (req, doc, note) in enumerate(docs):
        dy = Inches(1.6) + i * Inches(0.6)
        bc = GOLD if req == "必須" else MGRAY
        rect(s, Inches(0.35), dy, Inches(0.8), Inches(0.5), fill=bc)
        txb(s, req, Inches(0.35), dy+Inches(0.1), Inches(0.8), Inches(0.32),
            font=JP, size=9, color=DARK, bold=True, align=PP_ALIGN.CENTER)
        rect(s, Inches(1.2), dy, Inches(11.72), Inches(0.5),
             fill=RGBColor(0x22,0x20,0x1C))
        txb(s, f"□  {doc}", Inches(1.35), dy+Inches(0.07), Inches(6.0), Inches(0.38),
            font=JP, size=11.5, color=WHITE)
        txb(s, note, Inches(7.5), dy+Inches(0.1), Inches(5.2), Inches(0.32),
            font=JP, size=10, color=RGBColor(0x88,0x80,0x70))

    txb(s, "※ 時間単価の証憑（過去類似業務の請求書・契約書等）がBrain側で準備できるか事前確認が必要です。",
        Inches(0.35), Inches(7.1), Inches(12.62), Inches(0.3),
        font=JP, size=9, color=RGBColor(0x99,0x90,0x80))

    # ── S7: 今週のタスク ────────────────────────────────
    s = blank(prs)
    slide_bg(s, DARK2)
    left_accent(s, color=RED)
    label(s, "CHECK 06 — 今週のタスク一覧", Inches(0.4), Inches(0.4))
    heading(s, "誰が・何を・いつまでに", Inches(0.4), Inches(0.75), Inches(10))
    gold_bar(s, Inches(0.4), Inches(1.35), w=Inches(0.32), h=Inches(0.03))

    tasks = [
        ("小野氏", "プラン選択（松・竹・梅）",                "本日中",         True),
        ("小野氏", "申請者の確定（TOJISM or 丸正旅館）",       "本日中",         True),
        ("小野氏", "キャッシュフロー手当の確認",               "今週中",         True),
        ("小野氏", "補助金 参加申込（特設サイト）",             "5月22日（木）",  True),
        ("Brain",  "計画申請書（様式1〜4）の起草",             "5月28日（水）",  True),
        ("Brain",  "見積書・時間単価根拠資料の作成",            "5月28日（水）",  True),
        ("小野氏+Brain", "計画申請の最終確認・提出",           "5月29日（木）",  True),
        ("小野氏", "LPコンテンツ素材の提供（MISSION等）",      "5月末",          False),
    ]
    for i, (who, task, deadline, urgent) in enumerate(tasks):
        ty = Inches(1.6) + i * Inches(0.71)
        wc = GOLD if "小野氏" in who and "Brain" not in who else (GREEN if who == "Brain" else RGBColor(0x66,0xAA,0x88))
        rect(s, Inches(0.35), ty, Inches(1.5), Inches(0.6), fill=wc)
        txb(s, who, Inches(0.35), ty+Inches(0.13), Inches(1.5), Inches(0.38),
            font=JP, size=10, color=DARK, bold=True, align=PP_ALIGN.CENTER)
        rect(s, Inches(1.9), ty, Inches(10.7), Inches(0.6),
             fill=RGBColor(0x22,0x20,0x1C))
        txb(s, f"□  {task}", Inches(2.05), ty+Inches(0.13), Inches(7.8), Inches(0.38),
            font=JP, size=12, color=WHITE)
        dc = RED if urgent else MGRAY
        txb(s, deadline, Inches(10.0), ty+Inches(0.13), Inches(2.5), Inches(0.38),
            font=JP, size=11, color=dc, bold=urgent, align=PP_ALIGN.RIGHT)

    prs.save(r"C:\project\daikonya\補助金申請チェックシート.pptx")
    print("OK: hojo-check.pptx saved")


# ── Run ─────────────────────────────────────────────────
if __name__ == "__main__":
    make_proposal()
    make_subsidy()
    print("Done.")
