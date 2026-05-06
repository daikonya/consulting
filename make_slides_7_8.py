# -*- coding: utf-8 -*-
"""
スライド7・8のみを生成 → 追加スライド_7_8.pptx
既存の提案書_v2.pptx には一切触れない。
PowerPoint の「スライドの再利用」で手動挿入してください。
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colors ──────────────────────────────────────────────
DARK   = RGBColor(0x1C, 0x1A, 0x16)
DARK2  = RGBColor(0x26, 0x23, 0x1E)
DARK3  = RGBColor(0x22, 0x20, 0x1C)
GOLD   = RGBColor(0xB8, 0x93, 0x5A)
GOLD_L = RGBColor(0xD4, 0xAF, 0x7A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xEE, 0xEA, 0xE3)
MGRAY  = RGBColor(0x6B, 0x63, 0x59)
BGRAY  = RGBColor(0xDD, 0xD5, 0xC5)
RED    = RGBColor(0xC0, 0x39, 0x2B)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
TEAL   = RGBColor(0x1A, 0x7A, 0x5E)
AMBER  = RGBColor(0xE6, 0x8A, 0x1A)

SW = Inches(13.33)
SH = Inches(7.5)
JP  = "Yu Gothic"
JPS = "Yu Mincho"


def new_prs():
    p = Presentation()
    p.slide_width  = SW
    p.slide_height = SH
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=DARK):
    r = slide.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()

def box(slide, l, t, w, h, fill=None, lc=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if lc:   s.line.color.rgb = lc; s.line.width = lw
    else:    s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h,
       fn=JP, sz=13, fc=WHITE, bold=False,
       align=PP_ALIGN.LEFT, wrap=True):
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = fn; run.font.size = Pt(sz)
    run.font.color.rgb = fc; run.font.bold = bold
    return b

def sec_label(slide, text, color=GOLD):
    tb(slide, text, Inches(0.35), Inches(0.38), Inches(11), Inches(0.3),
       fn=JP, sz=9, fc=color)

def heading(slide, text, fc=WHITE, sz=26):
    tb(slide, text, Inches(0.35), Inches(0.7), Inches(11), Inches(0.9),
       fn=JPS, sz=sz, fc=fc)

def gold_line(slide):
    box(slide, Inches(0.35), Inches(1.38), Inches(0.32), Inches(0.025), fill=GOLD)

def accent_bar(slide, color=GOLD):
    box(slide, 0, 0, Inches(0.12), SH, fill=color)


# ════════════════════════════════════════════════════════
#  SLIDE 7 : 回収期間と導入価値の訴求
# ════════════════════════════════════════════════════════
def slide_07(prs):
    s = blank(prs)
    bg(s, DARK)
    accent_bar(s)
    sec_label(s, "PAYBACK -- 回収期間と導入価値（補助金活用時）")
    heading(s, "この投資は割に合うか？")
    gold_line(s)

    # ── 左：S6の収益サマリ ──────────────────────────────
    box(s, Inches(0.25), Inches(1.55), Inches(4.1), Inches(5.7), fill=DARK3)
    tb(s, "S6 収益シミュレーション（再掲）",
       Inches(0.45), Inches(1.68), Inches(3.7), Inches(0.32),
       fn=JP, sz=10, fc=GOLD)

    rows_left = [
        ("保守シナリオ（6件/年）", "69万円",  False),
        ("目標シナリオ（10件/年）","136万円", True),
        ("LTV込み 保守",           "約100万円", False),
        ("LTV込み 目標",           "約200万円", True),
    ]
    for i, (lbl, val, highlight) in enumerate(rows_left):
        ry = Inches(2.18) + i * Inches(1.18)
        fbg = RGBColor(0x1A,0x2E,0x1A) if highlight else DARK2
        box(s, Inches(0.45), ry, Inches(3.7), Inches(1.0), fill=fbg)
        tb(s, lbl, Inches(0.6), ry+Inches(0.1), Inches(3.4), Inches(0.35),
           fn=JP, sz=10, fc=MGRAY)
        fc_val = GREEN if highlight else GOLD_L
        tb(s, val, Inches(0.6), ry+Inches(0.48), Inches(3.4), Inches(0.42),
           fn=JPS, sz=22, fc=fc_val)

    # ── 中：矢印・仕切り ────────────────────────────────
    tb(s, "vs", Inches(4.45), Inches(3.8), Inches(0.5), Inches(0.5),
       fn=JP, sz=16, fc=MGRAY, align=PP_ALIGN.CENTER)

    # ── 右：回収期間分析（補助金あり）──────────────────
    box(s, Inches(5.1), Inches(1.55), Inches(7.98), Inches(5.7),
        fill=DARK3, lc=GOLD, lw=Pt(0.75))
    tb(s, "補助金活用時の実質コストと回収期間",
       Inches(5.3), Inches(1.68), Inches(7.5), Inches(0.32),
       fn=JP, sz=10, fc=GOLD)

    # 実質コスト
    box(s, Inches(5.3), Inches(2.1), Inches(7.58), Inches(1.08),
        fill=DARK2)
    tb(s, "年間実質自己負担（スタンダード以上・初年度）",
       Inches(5.5), Inches(2.18), Inches(7.2), Inches(0.3),
       fn=JP, sz=10, fc=MGRAY)
    tb(s, "SNS実費のみ  ＝  約20万円/年",
       Inches(5.5), Inches(2.5), Inches(5.0), Inches(0.42),
       fn=JP, sz=14, fc=WHITE, bold=True)
    box(s, Inches(10.1), Inches(2.18), Inches(2.7), Inches(0.88),
        fill=RGBColor(0x1A,0x2E,0x1A))
    tb(s, "LINE広告費60万は\nBrain負担（初年度）",
       Inches(10.2), Inches(2.22), Inches(2.55), Inches(0.72),
       fn=JP, sz=9.5, fc=TEAL)

    # 年次回収表
    years = [
        ("Year 1（保守）",  "収益 69万 ー コスト 20万",  "＋49万",        GREEN,  "初年度から黒字"),
        ("Year 1（目標）",  "収益136万 ー コスト 20万",  "＋116万",       GREEN,  "大幅黒字"),
        ("Year 2（LTV込）", "収益200万 ー コスト 72万※", "＋128万",       GREEN,  "安定黒字"),
        ("Year 3 以降",     "仕組みが自走 / コスト逓減", "加速度的に回収", GOLD_L, ""),
    ]
    for i, (yr, calc, result, rc, note) in enumerate(years):
        ry = Inches(3.18) + i * Inches(0.95)
        box(s, Inches(5.3), ry, Inches(7.58), Inches(0.86), fill=DARK2)
        tb(s, yr,     Inches(5.45), ry+Inches(0.08), Inches(2.2), Inches(0.32),
           fn=JP, sz=10.5, fc=MGRAY)
        tb(s, calc,   Inches(7.75), ry+Inches(0.08), Inches(3.0), Inches(0.32),
           fn=JP, sz=10, fc=MGRAY)
        tb(s, result, Inches(5.45), ry+Inches(0.44), Inches(3.5), Inches(0.36),
           fn=JPS, sz=16, fc=rc, bold=False)
        if note:
            tb(s, note, Inches(9.0), ry+Inches(0.5), Inches(3.7), Inches(0.28),
               fn=JP, sz=10, fc=rc)

    # ── 注釈 ────────────────────────────────────────────
    tb(s, "※Year2以降のコスト72万 ＝ LINE広告60万（自己負担）＋ SNS実費12万（機材費不要のため）",
       Inches(5.3), Inches(7.05), Inches(7.7), Inches(0.25),
       fn=JP, sz=8.5, fc=RGBColor(0x66,0x60,0x58))

    # ── 底部：結論 ──────────────────────────────────────
    box(s, Inches(0.25), Inches(7.32), Inches(12.83), Inches(0.5),
        fill=RGBColor(0x1A,0x2E,0x1A))
    tb(s, "結論：",
       Inches(0.45), Inches(7.43), Inches(1.0), Inches(0.3),
       fn=JP, sz=11, fc=TEAL, bold=True)
    tb(s, "初年度から黒字。LINE広告費をBrainが持つことで自己負担は20万のみ。導入しない理由はない。",
       Inches(1.5), Inches(7.43), Inches(11.4), Inches(0.3),
       fn=JP, sz=11, fc=WHITE)


# ════════════════════════════════════════════════════════
#  SLIDE 8 : コンサルフィー改定の提案
# ════════════════════════════════════════════════════════
def slide_08(prs):
    s = blank(prs)
    bg(s, DARK)
    accent_bar(s, color=AMBER)
    sec_label(s, "FEE STRATEGY -- コンサルフィー戦略", color=AMBER)
    heading(s, "フィー設定の現状と改定タイミングの提案")
    gold_line(s)

    # ── 現状分析 ────────────────────────────────────────
    box(s, Inches(0.25), Inches(1.55), Inches(6.1), Inches(2.55), fill=DARK3)
    tb(s, "現在のフィー（月換算）と市場相場の比較",
       Inches(0.45), Inches(1.68), Inches(5.7), Inches(0.32),
       fn=JP, sz=10, fc=AMBER)

    headers = ["プラン", "現在（月換算）", "市場相場", "乖離"]
    hx = [Inches(0.45), Inches(1.8), Inches(3.3), Inches(4.95)]
    hw = [Inches(1.3), Inches(1.45), Inches(1.6), Inches(1.0)]
    for j, h in enumerate(headers):
        box(s, hx[j], Inches(2.08), hw[j], Inches(0.35),
            fill=RGBColor(0x33,0x2E,0x22))
        tb(s, h, hx[j]+Inches(0.05), Inches(2.12), hw[j]-Inches(0.1), Inches(0.28),
           fn=JP, sz=9, fc=MGRAY)

    fee_rows = [
        ("エッセンシャル", "2.3万/月",  "10〜20万/月", "大"),
        ("スタンダード",   "5万/月",    "15〜30万/月", "大"),
        ("エグゼクティブ", "10万/月",   "20〜50万/月", "中"),
    ]
    for i, (plan, current, market, gap) in enumerate(fee_rows):
        ry = Inches(2.5) + i * Inches(0.5)
        box(s, Inches(0.45), ry, Inches(5.7), Inches(0.42), fill=DARK2)
        tb(s, plan,    hx[0]+Inches(0.05), ry+Inches(0.07), hw[0]-Inches(0.1), Inches(0.3),
           fn=JP, sz=10.5, fc=WHITE)
        tb(s, current, hx[1]+Inches(0.05), ry+Inches(0.07), hw[1]-Inches(0.1), Inches(0.3),
           fn=JP, sz=10.5, fc=AMBER)
        tb(s, market,  hx[2]+Inches(0.05), ry+Inches(0.07), hw[2]-Inches(0.1), Inches(0.3),
           fn=JP, sz=10.5, fc=MGRAY)
        tb(s, gap,     hx[3]+Inches(0.05), ry+Inches(0.07), hw[3]-Inches(0.1), Inches(0.3),
           fn=JP, sz=10.5, fc=RED, bold=True)

    tb(s, "現在の設定は市場参入価格として妥当。ただし中長期では継続性のためにも改定が必要。",
       Inches(0.45), Inches(4.18), Inches(5.7), Inches(0.42),
       fn=JP, sz=10, fc=MGRAY, wrap=True)

    # ── 改定タイミング ──────────────────────────────────
    box(s, Inches(6.55), Inches(1.55), Inches(6.53), Inches(2.55), fill=DARK3)
    tb(s, "フィー改定のトリガー",
       Inches(6.75), Inches(1.68), Inches(6.1), Inches(0.32),
       fn=JP, sz=10, fc=AMBER)

    box(s, Inches(6.75), Inches(2.08), Inches(6.13), Inches(1.8),
        fill=RGBColor(0x2E,0x22,0x0E), lc=AMBER, lw=Pt(1.2))
    tb(s, "月次の新規依頼が平均 3件 を超えたとき",
       Inches(6.95), Inches(2.2), Inches(5.7), Inches(0.42),
       fn=JPS, sz=16, fc=AMBER)
    tb(s, "需要が供給を超え始めるタイミング＝価格交渉力が最大になる瞬間。\nこのシグナルを見逃さず、次の契約更新時に改定を告知する。",
       Inches(6.95), Inches(2.7), Inches(5.7), Inches(0.7),
       fn=JP, sz=10.5, fc=RGBColor(0xCC,0xB8,0x88), wrap=True)

    # ── 2フェーズロードマップ ───────────────────────────
    box(s, Inches(0.25), Inches(4.72), Inches(12.83), Inches(0.4),
        fill=RGBColor(0x28,0x26,0x22))
    tb(s, "フィー改定ロードマップ",
       Inches(0.45), Inches(4.8), Inches(4), Inches(0.28),
       fn=JP, sz=10, fc=GOLD)

    phases = [
        ("Phase 1\n0〜12ヶ月",
         "市場参入フェーズ",
         "現在の料金で実績積み上げ\n口コミ・紹介・SNSでブランド確立\nケーススタディを蓄積",
         RGBColor(0x26,0x23,0x1E), MGRAY),
        ("Phase 2\n12ヶ月〜\n（3件超えが目安）",
         "フィー改定フェーズ",
         "エッセンシャル  7万  →  12万（＋71%）\nスタンダード  15万  →  25万（＋67%）\nエグゼクティブ  30万  →  50万（＋67%）",
         RGBColor(0x2A,0x22,0x0A), AMBER),
        ("Phase 3\n24ヶ月〜",
         "プレミアム化フェーズ",
         "事業承継・M&A支援など高単価領域へ\n成果報酬型の組み合わせも検討\n紹介報酬制度で自然増を設計",
         RGBColor(0x1A,0x2A,0x1A), TEAL),
    ]
    pw = Inches(4.2)
    for i, (phase, title, desc, fbg, ac) in enumerate(phases):
        px = Inches(0.25) + i * (pw + Inches(0.07))
        py = Inches(5.2)
        box(s, px, py, pw, Inches(2.15), fill=fbg, lc=ac, lw=Pt(0.75))
        box(s, px, py, pw, Inches(0.42), fill=ac)
        tb(s, phase, px+Inches(0.12), py+Inches(0.06),
           Inches(1.5), Inches(0.36), fn=JP, sz=9, fc=DARK, bold=True)
        tb(s, title, px+Inches(1.65), py+Inches(0.08),
           pw-Inches(1.75), Inches(0.3), fn=JP, sz=11, fc=DARK, bold=True)
        tb(s, desc, px+Inches(0.15), py+Inches(0.54),
           pw-Inches(0.28), Inches(1.5), fn=JP, sz=10.5, fc=RGBColor(0xCC,0xC5,0xB8), wrap=True)

    # ── 改定後の収益試算 ────────────────────────────────
    box(s, Inches(0.25), Inches(7.28), Inches(12.83), Inches(0.5),
        fill=RGBColor(0x2A,0x22,0x0A))
    tb(s, "改定後の試算：",
       Inches(0.45), Inches(7.38), Inches(1.6), Inches(0.3),
       fn=JP, sz=11, fc=AMBER, bold=True)
    tb(s, "目標10件 × 改定後単価（平均30万/契約）＝ 年間 約300万円  ／  補助金ありの実費80万と比較すると ROI 375%",
       Inches(2.1), Inches(7.38), Inches(11.0), Inches(0.3),
       fn=JP, sz=11, fc=WHITE)


# ── Main ────────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_07(prs)
    slide_08(prs)
    out = r"C:\project\daikonya\追加スライド_7_8_v2.pptx"
    prs.save(out)
    print("Done: 追加スライド_7_8.pptx saved (2 slides)")
    print("PowerPoint で提案書_v2.pptx を開き、")
    print("スライド6の後に「スライドの再利用」で挿入してください。")

if __name__ == "__main__":
    main()
