# -*- coding: utf-8 -*-
"""
サービス詳細サマリー + 締めスライド を別ファイルで生成
→ 追加スライド_サマリー締め.pptx
提案書_v2.pptx には触れない。PowerPointで手動挿入してください。
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK   = RGBColor(0x1C, 0x1A, 0x16)
DARK2  = RGBColor(0x26, 0x23, 0x1E)
DARK3  = RGBColor(0x22, 0x20, 0x1C)
GOLD   = RGBColor(0xB8, 0x93, 0x5A)
GOLD_L = RGBColor(0xD4, 0xAF, 0x7A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MGRAY  = RGBColor(0x6B, 0x63, 0x59)
TEAL   = RGBColor(0x1A, 0x7A, 0x5E)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
LINE_G = RGBColor(0x06, 0xC7, 0x55)
BLUE   = RGBColor(0x2E, 0x86, 0xC1)
PINK   = RGBColor(0xE1, 0x30, 0x6C)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
AMBER  = RGBColor(0xE6, 0x8A, 0x1A)

SW = Inches(13.33)
SH = Inches(7.5)
JP  = "Yu Gothic"
JPS = "Yu Mincho"


def new_prs():
    p = Presentation()
    p.slide_width  = SW
    p.slide_height = SH
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=DARK):
    r = slide.shapes.add_shape(1, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background()

def box(slide, l, t, w, h, fill=None, lc=None, lw=Pt(0)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if lc:   s.line.color.rgb = lc; s.line.width = lw
    else:    s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h,
       fn=JP, sz=13, fc=WHITE, bold=False,
       align=PP_ALIGN.LEFT, wrap=True):
    b = slide.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = fn; run.font.size = Pt(sz)
    run.font.color.rgb = fc; run.font.bold = bold
    return b


# ════════════════════════════════════════════════════════
#  SLIDE A : サービス詳細サマリー
# ════════════════════════════════════════════════════════
def slide_summary(prs):
    s = blank(prs)
    bg(s, DARK)
    # left accent
    box(s, 0, 0, Inches(0.12), SH, fill=GOLD)
    # label
    tb(s, "04 -- SERVICE DETAILS", Inches(0.35), Inches(0.38),
       Inches(8), Inches(0.3), fn=JP, sz=9, fc=GOLD)
    # heading
    tb(s, "サービス詳細 一覧", Inches(0.35), Inches(0.7),
       Inches(10), Inches(0.8), fn=JPS, sz=26, fc=WHITE)
    box(s, Inches(0.35), Inches(1.38), Inches(0.32), Inches(0.025), fill=GOLD)

    services = [
        {
            "no": "01", "title": "LP制作",
            "color": GOLD,
            "points": [
                "和モダン×高級感デザイン",
                "お問い合わせフォーム・CTA設計",
                "AI自動ブログ機能（プレミアム）",
            ],
            "plans": ["ライト", "スタンダード", "プレミアム"],
            "plan_color": GOLD,
        },
        {
            "no": "02", "title": "公式LINE構築＋広告",
            "color": LINE_G,
            "points": [
                "ステップ配信シナリオ 5〜8通",
                "無料相談予約フロー",
                "LINE広告運用（初年度Brain負担）",
            ],
            "plans": ["ライト", "スタンダード", "プレミアム"],
            "plan_color": LINE_G,
        },
        {
            "no": "03", "title": "ブランディング戦略支援",
            "color": AMBER,
            "points": [
                "ブランドコンセプト・キャッチコピー設計",
                "プロフィール・権威付けコンテンツ作成",
                "セミナー登壇プランニング",
            ],
            "plans": ["スタンダード", "プレミアム"],
            "plan_color": AMBER,
        },
        {
            "no": "04", "title": "月次マーケティング支援",
            "color": BLUE,
            "points": [
                "Google Analytics 4 / Search Console 分析",
                "Web広告効果測定・改善提案",
                "月次レポート＋翌月アクションプラン",
            ],
            "plans": ["スタンダード（4ヶ月）", "プレミアム（12ヶ月）"],
            "plan_color": BLUE,
        },
        {
            "no": "05", "title": "SNS動画制作・運用",
            "color": PINK,
            "points": [
                "月1回・60分の素材撮影",
                "月8本×1分 Shorts動画編集・投稿",
                "SNS運用マニュアル作成・引き渡し",
            ],
            "plans": ["スタンダード", "プレミアム"],
            "plan_color": PINK,
        },
        {
            "no": "06", "title": "AI活用・自動ブログ",
            "color": PURPLE,
            "points": [
                "GitHub Actions による自動記事生成",
                "Claude API 連携・週1〜2本更新",
                "SEO向け専門コンテンツを自動発信",
            ],
            "plans": ["プレミアム"],
            "plan_color": PURPLE,
        },
    ]

    cols = [Inches(0.25), Inches(4.57), Inches(8.89)]
    rows = [Inches(1.55), Inches(4.3)]
    cw = Inches(4.1)
    ch = Inches(2.55)

    for i, svc in enumerate(services):
        cx = cols[i % 3]
        cy = rows[i // 3]
        color = svc["color"]

        # card background
        box(s, cx, cy, cw, ch, fill=DARK3)
        # top color bar
        box(s, cx, cy, cw, Inches(0.38), fill=color)
        # number
        tb(s, svc["no"], cx+Inches(0.12), cy+Inches(0.06),
           Inches(0.5), Inches(0.28), fn=JP, sz=10, fc=DARK, bold=True)
        # title
        tb(s, svc["title"], cx+Inches(0.6), cy+Inches(0.06),
           cw-Inches(0.72), Inches(0.28), fn=JP, sz=11, fc=DARK, bold=True)

        # bullet points
        for j, pt in enumerate(svc["points"]):
            py = cy + Inches(0.52) + j * Inches(0.48)
            box(s, cx+Inches(0.18), py+Inches(0.16),
                Inches(0.06), Inches(0.06), fill=color)
            tb(s, pt, cx+Inches(0.32), py+Inches(0.06),
               cw-Inches(0.42), Inches(0.38),
               fn=JP, sz=10.5, fc=RGBColor(0xCC,0xC5,0xB8))

        # plan badges
        badge_x = cx + Inches(0.18)
        tb(s, "対象プラン：", badge_x, cy+ch-Inches(0.38),
           Inches(1.2), Inches(0.28), fn=JP, sz=8, fc=MGRAY)
        plan_str = "  /  ".join(svc["plans"])
        tb(s, plan_str, badge_x+Inches(1.1), cy+ch-Inches(0.4),
           cw-Inches(1.3), Inches(0.3), fn=JP, sz=9, fc=color, bold=True)


# ════════════════════════════════════════════════════════
#  SLIDE B : 締めスライド
# ════════════════════════════════════════════════════════
def slide_close(prs):
    s = blank(prs)
    bg(s, DARK)

    # 左ゴールドパネル
    box(s, 0, 0, Inches(4.5), SH, fill=DARK2)
    box(s, Inches(4.5), 0, Inches(0.06), SH, fill=GOLD)

    # 左：本日決めること
    tb(s, "本日決めること", Inches(0.35), Inches(1.2),
       Inches(3.8), Inches(0.4), fn=JP, sz=10, fc=GOLD)

    decisions = [
        ("01", "サポートプランの選択",      "ライト / スタンダード / プレミアム"),
        ("02", "補助金申請者の確定",         "丸正旅館 or 株式会社TOJISM"),
        ("03", "キャッシュフロー手当の確認", "先払い20〜80万の準備"),
        ("04", "参加申込（5/22〆）",          "今週中に特設サイトから申込"),
    ]
    for i, (no, title, sub) in enumerate(decisions):
        dy = Inches(1.75) + i * Inches(1.18)
        box(s, Inches(0.35), dy, Inches(3.8), Inches(1.05), fill=DARK3)
        box(s, Inches(0.35), dy, Inches(0.4), Inches(1.05),
            fill=GOLD)
        tb(s, no, Inches(0.35), dy+Inches(0.3), Inches(0.4), Inches(0.38),
           fn=JP, sz=11, fc=DARK, bold=True, align=PP_ALIGN.CENTER)
        tb(s, title, Inches(0.85), dy+Inches(0.1), Inches(3.2), Inches(0.38),
           fn=JP, sz=12, fc=WHITE, bold=False)
        tb(s, sub, Inches(0.85), dy+Inches(0.52), Inches(3.2), Inches(0.38),
           fn=JP, sz=10, fc=MGRAY)

    # 左：連絡先
    box(s, Inches(0.35), Inches(6.5), Inches(3.8), Inches(0.75), fill=DARK3)
    tb(s, "合同会社Brain  鵜木 芳文",
       Inches(0.5), Inches(6.58), Inches(3.5), Inches(0.3),
       fn=JP, sz=11, fc=WHITE)
    tb(s, "yoshifumi.unoki.business@gmail.com",
       Inches(0.5), Inches(6.88), Inches(3.5), Inches(0.28),
       fn=JP, sz=9, fc=MGRAY)

    # 右：メインメッセージ
    tb(s, "ご清聴\nありがとうございました。",
       Inches(4.8), Inches(1.5), Inches(8.2), Inches(2.5),
       fn=JPS, sz=36, fc=WHITE, wrap=True)

    # サブメッセージ
    tb(s, "旅館・飲食業の現場を変える伴走者として、\n小野氏のブランドを一緒に育てていきます。",
       Inches(4.8), Inches(4.1), Inches(8.2), Inches(0.9),
       fn=JP, sz=14, fc=RGBColor(0xAA,0xA2,0x95), wrap=True)

    # 区切り線
    box(s, Inches(4.8), Inches(5.1), Inches(8.0), Inches(0.02), fill=RGBColor(0x44,0x40,0x38))

    # ネクストステップ
    tb(s, "Next Step", Inches(4.8), Inches(5.25),
       Inches(3), Inches(0.3), fn=JP, sz=9, fc=GOLD)
    steps = [
        "本日：プラン選択・申請者確定",
        "今週中：補助金参加申込（5/22〆）",
        "5/29：計画申請書を共同提出",
        "交付決定後：LP・LINE・SNS 本格スタート",
    ]
    for i, step in enumerate(steps):
        sy = Inches(5.6) + i * Inches(0.44)
        box(s, Inches(4.8), sy+Inches(0.14), Inches(0.06), Inches(0.06), fill=GOLD)
        tb(s, step, Inches(5.0), sy+Inches(0.06), Inches(7.9), Inches(0.35),
           fn=JP, sz=11, fc=RGBColor(0xCC,0xC5,0xB8))

    # コピーライト
    tb(s, "© 2026 合同会社Brain", Inches(4.8), Inches(7.2),
       Inches(4), Inches(0.25), fn=JP, sz=9,
       fc=RGBColor(0x44,0x40,0x38))


# ── Main ────────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_summary(prs)
    slide_close(prs)
    out = r"C:\project\daikonya\追加スライド_サマリー締め.pptx"
    prs.save(out)
    print("Done: 追加スライド_サマリー締め.pptx saved (2 slides)")
    print("挿入先：")
    print("  サマリー → スライド10（LP別途投影）の直前")
    print("  締め     → 最終スライドの後")

if __name__ == "__main__":
    main()
